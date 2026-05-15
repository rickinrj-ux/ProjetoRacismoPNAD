"""
gerar_apresentacao_executiva.py
Versão para congressos e empresas — linguagem acessível, narrativa visual.
12 slides | Sem jargão técnico | Foco em impacto e políticas.
"""
import sys; sys.stdout.reconfigure(encoding='utf-8')
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT    = Path(r"C:\Users\user\Documents\ProjetoRacismoPNAD")
FIGURES = ROOT / "outputs" / "figures"
OUT_PPT = ROOT / "apresentacao_executiva.pptx"

# ── Paleta ────────────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1B, 0x2A, 0x4A)
C_BLUE   = RGBColor(0x15, 0x65, 0xC0)
C_RED    = RGBColor(0xC6, 0x28, 0x28)
C_AMBER  = RGBColor(0xFF, 0x8F, 0x00)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x21, 0x21, 0x21)
C_GRAY   = RGBColor(0x61, 0x61, 0x61)
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
C_TEAL   = RGBColor(0x00, 0x69, 0x6D)

W = In(13.33)
H = In(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, lpt=0):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.width = Pt(lpt) if lpt else 0
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if line and lpt: sh.line.color.rgb = line
    return sh

def text(slide, txt, l, t, w, h,
         size=18, bold=False, italic=False,
         color=C_BLACK, align=PP_ALIGN.LEFT, name="Calibri"):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = name
    return tf

def multiline(slide, lines, l, t, w, h,
              size=16, color=C_BLACK, bold=False,
              align=PP_ALIGN.LEFT, name="Calibri"):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = True
    first = True
    for line_txt in lines:
        p = tf.text_frame.paragraphs[0] if first else tf.text_frame.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = line_txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = name

def img(slide, path, l, t, w, h=None):
    p = Path(path)
    if not p.exists():
        text(slide, f"[{p.name}]", l, t, w, In(0.5), size=9, color=C_GRAY, italic=True)
        return
    slide.shapes.add_picture(str(p), l, t, w, h) if h else slide.shapes.add_picture(str(p), l, t, w)

def header(slide, title, subtitle=None, bg=C_DARK):
    rect(slide, 0, 0, W, In(1.05), fill=bg)
    text(slide, title, In(0.35), In(0.08), In(12.3), In(0.55),
         size=26, bold=True, color=C_WHITE, name="Calibri")
    if subtitle:
        text(slide, subtitle, In(0.35), In(0.63), In(12.3), In(0.35),
             size=13, color=RGBColor(0xBB, 0xDE, 0xFB), name="Calibri")

def footer_exec(slide, num, total=12):
    rect(slide, 0, H - In(0.25), W, In(0.25), fill=C_DARK)
    text(slide, "PNAD Contínua 2016–2025  |  15,9 milhões de observações  |  Ricardo Calheiros — MBA USP/ESALQ",
         In(0.2), H - In(0.23), In(11.5), In(0.22),
         size=7.5, color=RGBColor(0xBB, 0xDE, 0xFB), name="Calibri")
    text(slide, f"{num}/{total}",
         In(12.8), H - In(0.23), In(0.5), In(0.22),
         size=7.5, color=C_WHITE, align=PP_ALIGN.RIGHT, name="Calibri")

def big_stat(slide, value, label, sublabel, l, t, w=In(2.8), h=In(1.5),
             val_color=C_RED, bg=C_LGRAY):
    rect(slide, l, t, w, h, fill=bg, line=val_color, lpt=1.5)
    text(slide, value, l + In(0.1), t + In(0.08), w - In(0.2), In(0.72),
         size=34, bold=True, color=val_color, align=PP_ALIGN.CENTER, name="Calibri")
    text(slide, label, l + In(0.1), t + In(0.78), w - In(0.2), In(0.35),
         size=11.5, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER, name="Calibri")
    text(slide, sublabel, l + In(0.1), t + In(1.12), w - In(0.2), In(0.32),
         size=9.5, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True, name="Calibri")

def callout(slide, txt, l, t, w, h, bg=RGBColor(0xFF,0xF9,0xE7), border=C_AMBER):
    rect(slide, l, t, w, h, fill=bg, line=border, lpt=1.2)
    text(slide, txt, l + In(0.15), t + In(0.1), w - In(0.3), h - In(0.2),
         size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER, name="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=C_DARK)
rect(s, 0, 0, W, In(0.06), fill=C_RED)       # faixa vermelha no topo

text(s, "Racismo no Mercado de Trabalho Brasileiro",
     In(0.6), In(0.6), In(12.0), In(1.3),
     size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, name="Calibri")
text(s, "O que 15,9 milhões de registros da PNAD Contínua revelam sobre desigualdade, acesso e oportunidade",
     In(0.6), In(1.95), In(12.0), In(0.7),
     size=16, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER, name="Calibri")

rect(s, In(0.6), In(2.7), In(12.0), In(0.04), fill=C_AMBER)

for i, (v, l) in enumerate([
    ("37,5%",  "gap salarial médio\n(mediana: 27,5%)"),
    ("25%",    "menos chance de\nvaga qualificada"),
    ("+127%",  "ganho de renda com\ninclusão produtiva"),
    ("100+",   "anos para convergir\nno ritmo atual"),
]):
    x = In(0.6) + i * In(3.05)
    rect(s, x, In(2.85), In(2.85), In(1.35),
         fill=RGBColor(0x12, 0x1E, 0x38), line=C_AMBER, lpt=1)
    text(s, v, x + In(0.1), In(2.88), In(2.65), In(0.7),
         size=30, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    text(s, l, x + In(0.1), In(3.58), In(2.65), In(0.55),
         size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

text(s, "Ricardo Calheiros  |  MBA Data Science & Analytics  |  USP/ESALQ  |  2026",
     In(0.6), In(4.45), In(12.0), In(0.4),
     size=13, color=RGBColor(0x90, 0xA4, 0xAE), align=PP_ALIGN.CENTER)

text(s, "Dados: PNAD Contínua 2016–2025 (IBGE)  |  5 metodologias estatísticas  |  Período: 10 anos",
     In(0.6), H - In(0.55), In(12.0), In(0.35),
     size=11, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O PROBLEMA EM TRÊS REALIDADES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "O Problema",
       "Três realidades que os números confirmam — e que políticas isoladas não resolvem")

text(s, "A desigualdade racial no trabalho não é apenas salário:",
     In(0.4), In(1.2), In(12.5), In(0.4),
     size=17, bold=True, color=C_DARK)

cards = [
    (C_RED,   "1. Acesso bloqueado",
     "Negros têm 25% menos chance de estar em cargos qualificados mesmo com a mesma formação, a mesma experiência e morando no mesmo bairro que um colega branco.",
     "Mulheres negras sofrem dupla barreira: discriminação racial + de gênero simultaneamente (Crenshaw, 1989)."),
    (C_BLUE,  "2. Salário diferente\npela mesma função",
     "Após controlar todas as características observáveis, negros ainda recebem 5,4% menos na mesma função. No topo salarial esse gap chega a 12%.",
     "A discriminação não para na contratação — continua na progressão e no salário."),
    (C_TEAL,  "3. Onde você mora\ndetermina o que você ganha",
     "Mais de 1 em cada 5 reais da diferença salarial entre trabalhadores vem do bairro (UPA) onde cada um vive — não da competência individual.",
     "Segregação residencial se converte diretamente em segregação de renda."),
]
for i, (color, title, body, conclusion) in enumerate(cards):
    x = In(0.3) + i * In(4.35)
    rect(s, x, In(1.7), In(4.1), In(5.35),
         fill=C_LGRAY, line=color, lpt=2)
    rect(s, x, In(1.7), In(4.1), In(0.6), fill=color)
    text(s, title, x + In(0.12), In(1.73), In(3.86), In(0.55),
         size=13, bold=True, color=C_WHITE, name="Calibri")
    text(s, body, x + In(0.15), In(2.38), In(3.8), In(2.3),
         size=12.5, color=C_BLACK, name="Calibri")
    rect(s, x + In(0.12), In(4.75), In(3.86), In(0.03), fill=color)
    text(s, conclusion, x + In(0.15), In(4.82), In(3.8), In(0.8),
         size=11.5, bold=True, color=color, italic=True, name="Calibri")

footer_exec(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ONDE VOCÊ MORA DETERMINA SUA RENDA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Onde Você Mora Determina Sua Renda",
       "A segregação residencial se converte diretamente em segregação de renda")

img(s, FIGURES / "hlm_justif_caterpillar.png", In(0.3), In(1.2), In(5.8), In(4.6))

text(s, "O que isso significa?", In(7.5), In(1.2), In(5.6), In(0.4),
     size=16, bold=True, color=C_DARK)

bullets = [
    "Dois trabalhadores idênticos — mesma idade, mesmo diploma, mesma função — ganham salários diferentes porque nasceram em bairros diferentes.",
    "A renda média do bairro (UPA) é o MAIOR determinante individual de renda — mais do que a educação.",
    "Negros são super-representados nas UPAs de menor renda: o racismo histórico na moradia amplia o racismo atual no trabalho.",
    "Políticas de emprego que ignoram a dimensão territorial atacam apenas parte do problema.",
]
for i, b in enumerate(bullets):
    rect(s, In(7.5), In(1.72) + i * In(1.1), In(5.6), In(0.98),
         fill=C_LGRAY if i % 2 == 0 else RGBColor(0xE8,0xF5,0xE9),
         line=C_TEAL if i % 2 == 0 else C_GREEN, lpt=0.6)
    text(s, "▸  " + b, In(7.65), In(1.77) + i * In(1.1), In(5.3), In(0.88),
         size=12, color=C_BLACK, name="Calibri")

callout(s, "Em linguagem simples: se você mudar um negro do bairro mais pobre para o mais rico, "
           "você explica mais da metade do gap salarial — sem mudar nenhuma característica individual.",
        In(0.3), In(6.08), In(12.7), In(0.62))
footer_exec(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — A PIRÂMIDE RACIAL DO TRABALHO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "A Pirâmide Racial do Trabalho Brasileiro",
       "Nas melhores vagas: poucos negros. Nos piores empregos: maioria negra")

img(s, FIGURES / "comp_razao_grupo_cbo.png",       In(0.3),  In(1.2), In(6.3))
img(s, FIGURES / "comp_representacao_topo.png",    In(6.8),  In(1.2), In(6.2))

for i, (v, lbl, color) in enumerate([
    ("42",    "negros por 100 brancos\nna mesma função de dirigente", C_RED),
    ("2,1×",  "mais negros em\nocupações elementares",                C_AMBER),
    ("0,47",  "índice de representação\nnegra no top 5% das capitais", C_RED),
]):
    x = In(0.3) + i * In(4.1)
    big_stat(s, v, lbl, "", x, In(5.05), w=In(3.9), h=In(1.4),
             val_color=color, bg=RGBColor(0xFF,0xEB,0xEE) if color == C_RED else RGBColor(0xFF,0xF9,0xE7))

callout(s,
        "A sub-representação não é resultado de falta de qualificação — é o produto de barreiras "
        "sistemáticas de acesso confirmadas pelos modelos estatísticos.",
        In(0.3), In(6.55), In(12.7), In(0.60), bg=RGBColor(0xFF,0xEB,0xEE), border=C_RED)
footer_exec(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — NEGROS TÊM MENOS CHANCE NAS MELHORES VAGAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "Negros Têm Menos Chance nas Melhores Vagas",
       "Mesmo com o mesmo diploma, o mesmo bairro e a mesma experiência — a barreira persiste")

img(s, FIGURES / "glmm_odds_ratios_full.png", In(0.3), In(1.15), In(6.8))

text(s, "O que a análise mostrou:", In(7.6), In(1.2), In(5.5), In(0.4),
     size=16, bold=True, color=C_DARK)

achados = [
    (C_RED,   "25% menos chance",
     "Um trabalhador negro com exatamente as mesmas características que um branco tem 25% menos chance de estar em uma ocupação qualificada (gerente, engenheiro, advogado...)"),
    (C_AMBER, "−1,07 pontos percentuais",
     "Mesmo após considerar o bairro onde mora, o setor de trabalho e horas trabalhadas, negros ainda ficam fora das melhores vagas. Isso é o resíduo puro da discriminação."),
    (C_BLUE,  "Não é falta de diploma",
     "O modelo controla educação, experiência, setor, horas e bairro. O que sobra é a cor da pele sendo usada como critério de seleção."),
]
for i, (color, title, body) in enumerate(achados):
    rect(s, In(7.6), In(1.75) + i * In(1.45), In(5.5), In(1.32),
         fill=C_LGRAY, line=color, lpt=1.5)
    text(s, title, In(7.75), In(1.8) + i * In(1.45), In(5.2), In(0.4),
         size=15, bold=True, color=color, name="Calibri")
    text(s, body, In(7.75), In(2.22) + i * In(1.45), In(5.2), In(0.75),
         size=11.5, color=C_BLACK, name="Calibri")

footer_exec(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — O TETO DE VIDRO É REAL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "O Teto de Vidro É Real — e Piora no Topo",
       "Quanto maior o salário, maior a barreira racial — confirmado por análise estatística formal")

img(s, FIGURES / "quantreg_trajetoria.png", In(0.3), In(1.2), In(7.2))

text(s, "Como ler este gráfico:", In(7.8), In(1.2), In(5.3), In(0.4),
     size=16, bold=True, color=C_DARK)

bullets = [
    ("Eixo vertical",    "Quanto negros ganham a menos (%) em cada faixa salarial"),
    ("Linha descendente","A desvantagem CRESCE conforme o salário aumenta"),
    ("No salário mínimo","Gap racial de ~10% — ruim, mas menor"),
    ("No topo 5%",       "Gap racial sobe para ~22% — a barreira é mais alta exatamente onde os prêmios são maiores"),
    ("A cor da pele",    "age como um filtro que fica mais opaco conforme a vaga fica mais valiosa"),
]
for i, (titulo, corpo) in enumerate(bullets):
    rect(s, In(7.8), In(1.7) + i * In(0.9), In(5.3), In(0.82),
         fill=RGBColor(0xFF,0xEB,0xEE) if i >= 3 else C_LGRAY,
         line=C_RED if i >= 3 else C_GRAY, lpt=0.6)
    text(s, titulo + ": ", In(7.95), In(1.74) + i * In(0.9), In(1.5), In(0.38),
         size=11, bold=True, color=C_RED if i >= 3 else C_DARK)
    text(s, corpo, In(9.5), In(1.74) + i * In(0.9), In(3.5), In(0.38),
         size=11, color=C_BLACK, name="Calibri")

callout(s,
        "Ciclo de vida: o gap cresce de 9% (14–24 anos) para 37,5% (35–44 anos, pico de carreira) "
        "— o teto de vidro se fecha exatamente quando as promoções mais importam.",
        In(0.3), In(6.12), In(12.7), In(0.62))
footer_exec(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — A REDE QUE EXCLUI
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "A Rede que Exclui — Capital Social e Mobilidade",
       "Negros ficam fora das redes que convertem diplomas em empregos")

img(s, FIGURES / "sna_rede_demografica.png", In(0.3), In(1.2), In(6.2))

text(s, "O que são as redes profissionais?", In(6.9), In(1.2), In(6.1), In(0.4),
     size=16, bold=True, color=C_DARK)

bullets = [
    "Vagas nunca anunciadas são preenchidas por indicação — quem está na rede chega antes.",
    "Mentores, referências e \"network\" determinam promoção tão quanto o desempenho.",
    "A análise mostra: brancos atuam como conectores (brokers) em TODOS os níveis educacionais.",
    "Negros ficam nas bordas da rede — mesmo com pós-graduação, sem acesso ao centro.",
    "Resultado: o diploma negro vale menos porque falta a rede que o converte em oportunidade.",
]
for i, b in enumerate(bullets):
    ico = "◆" if i < 2 else ("✗" if i == 3 else "▸")
    c = C_RED if i == 3 else (C_BLUE if i < 2 else C_BLACK)
    rect(s, In(6.9), In(1.72) + i * In(0.95), In(6.1), In(0.85),
         fill=RGBColor(0xFF,0xEB,0xEE) if i == 3 else C_LGRAY,
         line=C_RED if i == 3 else C_GRAY, lpt=0.5)
    text(s, f"{ico}  {b}", In(7.05), In(1.77) + i * In(0.95), In(5.8), In(0.75),
         size=12.5, color=c, name="Calibri")

callout(s,
        '"Mesmo que dois candidatos tenham o mesmo currículo, o que tem a rede certa chega primeiro. '
        'Essa rede, historicamente, é majoritariamente branca."',
        In(0.3), In(6.55), In(6.5), In(0.68), bg=RGBColor(0xE3,0xF2,0xFD), border=C_BLUE)
footer_exec(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — O ESTADO AJUDA, MAS NÃO RESOLVE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "O Estado Ajuda — Mas Não Resolve",
       "Setor público reduz a desigualdade geral, mas o gap racial persiste e o de gênero é MAIOR")

img(s, FIGURES / "estado_h1_gini.png",   In(0.3), In(1.2), In(6.0))
img(s, FIGURES / "estado_h2h3_gaps.png", In(6.6), In(1.2), In(6.4))

for i, (v, lbl, sub, color) in enumerate([
    ("Sim",   "Gini público < privado",           "0,466 vs 0,472 — melhora marginal", C_GREEN),
    ("Não",   "Gap racial não some no concurso",  "−25,2% público vs −29,2% privado",  C_AMBER),
    ("Piora", "Gap de gênero é maior no público", "−20,8% público vs −16,9% privado",  C_RED),
]):
    x = In(0.3) + i * In(4.15)
    rect(s, x, In(5.05), In(3.9), In(1.2),
         fill=RGBColor(0xE8,0xF5,0xE9) if color == C_GREEN else
              (RGBColor(0xFF,0xF9,0xE7) if color == C_AMBER else RGBColor(0xFF,0xEB,0xEE)),
         line=color, lpt=1.5)
    text(s, v, x + In(0.15), In(5.08), In(1.3), In(0.55),
         size=26, bold=True, color=color, name="Calibri")
    text(s, lbl, x + In(0.15), In(5.62), In(3.6), In(0.35),
         size=12, bold=True, color=C_BLACK, name="Calibri")
    text(s, sub, x + In(0.15), In(5.97), In(3.6), In(0.25),
         size=10, color=C_GRAY, italic=True, name="Calibri")

callout(s,
        "O concurso público equaliza a entrada, mas a promoção às posições de liderança (DAS, chefia) "
        "ainda reflete os vieses do setor privado — ou piores.",
        In(0.3), In(6.38), In(12.7), In(0.55))
footer_exec(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — O CUSTO PARA O BRASIL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "O Custo para Todos Nós — Inclusão Produtiva como Oportunidade",
       "Hsieh et al. (2019): discriminação racial e de gênero custou ao EUA 20–40% do crescimento 1960–2010")

img(s, FIGURES / "estado_h4_inclusao.png", In(0.3), In(1.2), In(6.5))

text(s, "O que aconteceria se negros tivessem\nacesso igual às melhores vagas?",
     In(7.3), In(1.2), In(5.8), In(0.65),
     size=16, bold=True, color=C_DARK)

cenarios = [
    (C_RED,   "Hoje",
     "Mediana salarial: R$2.000 (brancos) vs R$1.450 (negros) — gap de 27,5%",
     "Gap médio de 37,5%; em cargos de alta qualificação: 11% negros vs 22% brancos"),
    (C_GREEN, "+127% renda",
     "Se negros tivessem a mesma distribuição ocupacional que brancos",
     "A renda média dos trabalhadores negros mais que dobraria — sem nenhum aumento de produtividade"),
    (C_TEAL,  "Proxy: +74 p.p. PIB",
     "Estimativa conservadora do impacto no crescimento",
     "Incluir negros nas melhores vagas seria um dos maiores motores de crescimento disponíveis hoje"),
]
for i, (color, title, body1, body2) in enumerate(cenarios):
    rect(s, In(7.3), In(2.0) + i * In(1.5), In(5.8), In(1.38),
         fill=C_LGRAY, line=color, lpt=2)
    rect(s, In(7.3), In(2.0) + i * In(1.5), In(5.8), In(0.4), fill=color)
    text(s, title, In(7.45), In(2.04) + i * In(1.5), In(5.5), In(0.38),
         size=14, bold=True, color=C_WHITE, name="Calibri")
    text(s, body1, In(7.45), In(2.45) + i * In(1.5), In(5.5), In(0.38),
         size=12, bold=True, color=color, name="Calibri")
    text(s, body2, In(7.45), In(2.82) + i * In(1.5), In(5.5), In(0.38),
         size=11, color=C_GRAY, italic=True, name="Calibri")

footer_exec(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — O QUE FUNCIONA: TRÊS EIXOS DE AÇÃO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "O Que Funciona — Três Eixos de Ação",
       "Nenhum eixo isolado resolve. Os três precisam atuar simultaneamente.")

eixos = [
    (C_RED,   "Abrir as Portas",
     "Cotas de acesso a cargos\nqualificados (meta: IR ≥ 0,80)",
     "Formação profissional com\nrecorte racial e bolsas de\nresidência em empresas premium",
     "Atacar onde os dados apontam:\nbarreiras de ACESSO, não só salário"),
    (C_BLUE,  "Pagar Igual",
     "Transparência salarial por raça\nobrigatória (empresas > 100 fun.)",
     "Auditoria de igual pagamento\npor trabalho igual com\npenalidades progressivas",
     "O gap residual de 5,4% é real\ne mensurável — e deve ser eliminado"),
    (C_GREEN, "Ampliar as Redes",
     "Mentoria estruturada para elevar\na posição de negros nas redes\nprofissionais",
     "30% dos cargos de liderança\n(DAS + corporativo) para negros\naté 2030",
     "Sem acesso às redes que convertem\ndiplomas em empregos, nada muda"),
]
for i, (color, title, a1, a2, impacto) in enumerate(eixos):
    x = In(0.3) + i * In(4.35)
    rect(s, x, In(1.2), In(4.1), In(5.85),
         fill=C_LGRAY, line=color, lpt=2)
    rect(s, x, In(1.2), In(4.1), In(0.55), fill=color)
    text(s, title, x + In(0.12), In(1.23), In(3.86), In(0.5),
         size=15, bold=True, color=C_WHITE, name="Calibri")
    text(s, "Ação 1:\n" + a1, x + In(0.2), In(1.85), In(3.7), In(1.2),
         size=12, color=C_BLACK, name="Calibri")
    rect(s, x + In(0.2), In(3.1), In(3.7), In(0.02), fill=color)
    text(s, "Ação 2:\n" + a2, x + In(0.2), In(3.2), In(3.7), In(1.2),
         size=12, color=C_BLACK, name="Calibri")
    rect(s, x + In(0.1), In(4.5), In(3.9), In(0.02), fill=color)
    text(s, impacto, x + In(0.15), In(4.6), In(3.8), In(0.7),
         size=11, bold=True, color=color, italic=True, name="Calibri")

text(s, "⚠  Sem ação simultânea nos três eixos, fechar o gap levará mais de 100 anos no ritmo atual.",
     In(0.3), H - In(0.55), In(12.7), In(0.38),
     size=13, bold=True, color=C_RED, align=PP_ALIGN.CENTER, name="Calibri")
footer_exec(s, 10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — A URGÊNCIA: OS NÚMEROS DA MUDANÇA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "A Urgência — O Que Está em Jogo",
       "Pleno emprego sem prosperidade: 94,4% empregados, renda real estagnada")

img(s, FIGURES / "estado_h4_tendencia.png",  In(0.3), In(1.2), In(6.2))
img(s, FIGURES / "estado_h5_armadilha.png",  In(6.7), In(1.2), In(6.3))

for i, (v, l, sub, color) in enumerate([
    ("94,4%",  "emprego em 2025", "máxima histórica da série",    C_GREEN),
    ("R$1.283","renda real 2025",  "queda vs R$1.345 em 2024",     C_RED),
    ("11%",    "negros em alta\nqualificação", "brancos: 22,4%",   C_RED),
    ("+100",   "anos para fechar\no gap",       "no ritmo atual",   C_AMBER),
]):
    x = In(0.3) + i * In(3.2)
    big_stat(s, v, l, sub, x, In(5.2), w=In(3.0), h=In(1.35),
             val_color=color,
             bg=RGBColor(0xE8,0xF5,0xE9) if color == C_GREEN else
                (RGBColor(0xFF,0xF9,0xE7) if color == C_AMBER else RGBColor(0xFF,0xEB,0xEE)))

callout(s,
        "O Brasil está a produzir emprego, mas não prosperidade igualitária. Pleno emprego com "
        "segregação ocupacional é o cenário atual — e não se resolve sozinho.",
        In(0.3), H - In(0.6), In(12.7), In(0.44))
footer_exec(s, 11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONTRACAPA / CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=C_DARK)
rect(s, 0, 0, W, In(0.06), fill=C_RED)

text(s, "A pergunta não é se o racismo estrutural existe.",
     In(0.8), In(0.7), In(11.7), In(0.7),
     size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, name="Calibri")
text(s, "Os dados respondem isso com clareza.",
     In(0.8), In(1.35), In(11.7), In(0.55),
     size=20, color=C_AMBER, align=PP_ALIGN.CENTER, name="Calibri")

rect(s, In(0.8), In(2.0), In(11.7), In(0.04), fill=C_AMBER)

text(s, "A pergunta é: o que faremos com esse conhecimento?",
     In(0.8), In(2.15), In(11.7), In(0.65),
     size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, name="Calibri")

for i, (v, l) in enumerate([
    ("25%",   "menos chance de uma\nvaga qualificada"),
    ("27,5%", "gap salarial na\nmediana bruta"),
    ("+127%", "ganho potencial de\nrenda com inclusão"),
    ("+100",  "anos para convergir\nsem ação deliberada"),
]):
    x = In(0.6) + i * In(3.05)
    rect(s, x, In(2.95), In(2.85), In(1.45),
         fill=RGBColor(0x12, 0x1E, 0x38), line=C_AMBER, lpt=1)
    text(s, v, x + In(0.1), In(2.98), In(2.65), In(0.68),
         size=30, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    text(s, l, x + In(0.1), In(3.66), In(2.65), In(0.62),
         size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

text(s, "Ricardo Calheiros  |  MBA Data Science & Analytics  |  USP/ESALQ  |  rickinrj@gmail.com",
     In(0.6), In(4.6), In(11.7), In(0.4),
     size=13, color=RGBColor(0x90, 0xA4, 0xAE), align=PP_ALIGN.CENTER)

text(s, "Dados e código disponíveis  |  Metodologia: HLM + ML/SHAP + Oaxaca-Blinder + GLMM + Reg. Quantílica  |  PNAD Contínua 2016–2025 (IBGE)",
     In(0.3), H - In(0.55), In(12.7), In(0.35),
     size=9.5, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.CENTER, italic=True)
text(s, "12/12", In(12.8), H - In(0.25), In(0.5), In(0.22),
     size=7.5, color=RGBColor(0x78, 0x90, 0x9C), align=PP_ALIGN.RIGHT)

# ── Salvar ────────────────────────────────────────────────────────────────────
prs.save(str(OUT_PPT))
print(f"Arquivo gerado: {OUT_PPT.name}")
print(f"Tamanho: {OUT_PPT.stat().st_size // 1024} KB")
print("Versão executiva (12 slides) para congressos e empresas.")
