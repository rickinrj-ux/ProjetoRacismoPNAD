"""
gerar_apresentacao_executiva_pptx.py
Gera apresentacao_executiva_tcc.pptx — 10 slides para gestores, banca ou stakeholders.
Foco: impacto visual, números grandes, linguagem acessível, sem jargão técnico.
"""
import sys; sys.stdout.reconfigure(encoding='utf-8')
import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
from params import P, fmt, fmtN, ame, or_str

ROOT    = Path(r"C:\Users\user\Documents\ProjetoRacismoPNAD")
FIGURES = ROOT / "outputs" / "figures"
TABLES  = ROOT / "outputs" / "tables"
OUT_PPT = ROOT / "apresentacao_executiva_tcc.pptx"

# ── Paleta ────────────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1F, 0x38, 0x64)
C_BLUE   = RGBColor(0x15, 0x65, 0xC0)
C_RED    = RGBColor(0xB7, 0x1C, 0x1C)
C_AMBER  = RGBColor(0xFF, 0x8F, 0x00)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x21, 0x21, 0x21)
C_GRAY   = RGBColor(0x61, 0x61, 0x61)
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
C_PURPLE = RGBColor(0x6A, 0x00, 0x8A)

W = In(13.33)
H = In(7.50)
TOTAL_SLIDES = 10

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = Pt(line_pt) if line_pt else 0
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=20, bold=False, italic=False,
             color=C_BLACK, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tf

def add_img(slide, img_path, l, t, w, h=None):
    p = Path(img_path)
    if not p.exists():
        add_text(slide, f"[{p.name}]", l, t, w, In(0.8),
                 font_size=9, color=C_GRAY, italic=True)
        return
    if h:
        slide.shapes.add_picture(str(p), l, t, w, h)
    else:
        slide.shapes.add_picture(str(p), l, t, w)

def header_bar(slide, title, subtitle=None, bg=C_DARK):
    add_rect(slide, 0, 0, W, In(1.05), fill_rgb=bg)
    add_text(slide, title, In(0.35), In(0.07), In(12.5), In(0.55),
             font_size=27, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, In(0.35), In(0.63), In(12.5), In(0.35),
                 font_size=13.5, color=RGBColor(0xBB,0xDE,0xFB))

def footer(slide, n):
    add_rect(slide, 0, H-In(0.28), W, In(0.28), fill_rgb=C_DARK)
    add_text(slide,
             "Ricardo Calheiros  |  MBA USP/ESALQ  |  Racismo Estrutural e Mercado de Trabalho  |  Versao Executiva",
             In(0.15), H-In(0.25), In(11.5), In(0.25),
             font_size=8, color=RGBColor(0xBB,0xDE,0xFB))
    add_text(slide, f"{n}/{TOTAL_SLIDES}", In(12.7), H-In(0.25), In(0.5), In(0.25),
             font_size=8, color=C_WHITE, align=PP_ALIGN.RIGHT)

def big_kpi(slide, label, value, sub, l, t, w=In(3.0), h=In(1.8),
            val_color=C_AMBER, bg=RGBColor(0x0D,0x1F,0x3C), border=C_AMBER):
    add_rect(slide, l, t, w, h, fill_rgb=bg, line_rgb=border, line_pt=1.5)
    add_text(slide, label, l+In(0.12), t+In(0.08), w-In(0.24), In(0.38),
             font_size=12, color=RGBColor(0xBB,0xDE,0xFB))
    add_text(slide, value, l+In(0.12), t+In(0.42), w-In(0.24), In(0.85),
             font_size=36, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, sub, l+In(0.12), t+In(1.3), w-In(0.24), In(0.38),
             font_size=10, color=C_GRAY, italic=True)

def bullet_box(slide, items, l, t, w, h, dot="▸", font_size=16,
               color=C_BLACK, dot_color=C_BLUE):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.text_frame.paragraphs[0]; first = False
        else:
            p = tf.text_frame.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run(); r1.text = dot + "  "
        r1.font.size = Pt(font_size); r1.font.color.rgb = dot_color
        r1.font.bold = True; r1.font.name = "Calibri"
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(font_size); r2.font.color.rgb = color
        r2.font.name = "Calibri"

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA EXECUTIVA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=C_DARK)
add_rect(s, In(0.5), In(0.5), In(12.33), In(4.1),
         fill_rgb=RGBColor(0x0D,0x1F,0x3C), line_rgb=C_AMBER, line_pt=2)

add_text(s, "Racismo Estrutural e Mercado de Trabalho no Brasil",
         In(0.9), In(0.75), In(11.5), In(1.6),
         font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Evidencias quantitativas da PNAD Continua 2016-2025",
         In(0.9), In(2.4), In(11.5), In(0.55),
         font_size=18, color=RGBColor(0xBB,0xDE,0xFB), align=PP_ALIGN.CENTER)
add_rect(s, In(0.9), In(3.05), In(11.5), In(0.04), fill_rgb=C_AMBER)
add_text(s, "Ricardo Calheiros  |  MBA Data Science & Analytics  |  USP/ESALQ  |  2026",
         In(0.9), In(3.2), In(11.5), In(0.4),
         font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

# 4 KPIs na parte inferior da capa
for i, (val, lbl) in enumerate([
    ("7,7M",  "trabalhadores analisados"),
    ("10",    "anos de dados (2016-2025)"),
    ("5",     "metodos complementares"),
    ("3",     "barreiras estruturais mapeadas"),
]):
    x = In(0.5) + i * In(3.2)
    add_rect(s, x, In(5.0), In(2.9), In(1.2),
             fill_rgb=RGBColor(0x1A,0x30,0x5C), line_rgb=C_AMBER, line_pt=0.8)
    add_text(s, val, x+In(0.1), In(5.05), In(2.7), In(0.6),
             font_size=30, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x+In(0.1), In(5.65), In(2.7), In(0.4),
             font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(s, "VERSAO EXECUTIVA — 10 SLIDES",
         In(0.3), H-In(0.5), In(12.7), In(0.4),
         font_size=11, color=RGBColor(0x90,0xA4,0xAE), align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O CUSTO DO RACISMO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "1. O Custo do Racismo para o Brasil",
           "Uma questao economica, nao apenas moral — trabalhadores negros representam 56% da PEA")

add_text(s, "Por que isso importa para o pais?",
         In(0.4), In(1.15), In(12.5), In(0.45),
         font_size=20, bold=True, color=C_DARK)

fatos = [
    (C_RED,   "+127,3%",
     "ganho de renda para negros se tivessem a mesma distribuicao ocupacional (CBO) que brancos\n— simulacao com dados reais PNAD 2025 (run_hipoteses_estado.py)"),
    (C_AMBER, "+74,3 p.p.",
     "proxy de ganho produtivo agregado (127,3% x participacao negra na forca de trabalho ~58%)\n— proxy conservador, nao e estimativa de PIB formal (Hsieh et al., 2019 como referencia)"),
    (C_BLUE,  "+100 anos",
     "tempo para eliminar o gap ao ritmo atual de convergencia espontanea\n— tendencia temporal: delta = 0,001 log-pt/ano (WLS 2016-2025, p=0,077 — nao significativo)"),
]
for i, (color, val, desc) in enumerate(fatos):
    x = In(0.3) + i * In(4.35)
    add_rect(s, x, In(1.75), In(4.1), In(3.8),
             fill_rgb=RGBColor(0xF5,0xF5,0xF5), line_rgb=color, line_pt=2)
    add_rect(s, x, In(1.75), In(4.1), In(0.55), fill_rgb=color)
    add_text(s, val, x+In(0.15), In(1.8), In(3.8), In(0.48),
             font_size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+In(0.15), In(2.4), In(3.8), In(3.0),
             font_size=13.5, color=C_BLACK)

add_rect(s, In(0.3), In(5.7), In(12.7), In(0.6),
         fill_rgb=RGBColor(0x1F,0x38,0x64))
add_text(s,
    "Este trabalho quantifica os mecanismos — onde a discriminacao opera, qual a magnitude e o que pode ser feito.",
    In(0.5), In(5.75), In(12.3), In(0.5),
    font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
footer(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — AS 3 BARREIRAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "2. Tres Barreiras Estruturais — O Diagnostico",
           "5 metodos convergem para o mesmo mapa de exclusao racial")

add_text(s, "A desigualdade racial no mercado de trabalho opera em tres estagios independentes:",
         In(0.4), In(1.12), In(12.5), In(0.42),
         font_size=16, bold=False, color=C_DARK)

barreiras = [
    (C_RED,    "BARREIRA I",   "Acesso e Segregacao",
     [f"Negros tem {fmt(P['OR_M2_menor_pct'],1)}% menos chance de obter\nocupacao qualificada (GLMM M2)",
      "Mesmo com educacao, sexo e bairro IDENTICOS",
      "Causa: segregacao residencial e\nexclusao das redes profissionais"]),
    (C_AMBER,  "BARREIRA II",  "Penalidade Salarial",
     ["Gap salarial residual de 6,2% dentro\nda mesma funcao e empresa",
      "Cresce para 15-18% para mulheres negras\n(interseccionalidade de Crenshaw, 1989)",
      "Gap varia entre estados: SD = 5,0 log-pt\n(random slope HLM — N=7,7M)"]),
    (C_DARK,   "BARREIRA III", "Isolamento Estrutural",
     ["Betweenness centrality = 0 para negros\nem TODAS as faixas de escolaridade (SNA)",
      "Diplomas negros valem menos porque\nfaltam redes para converte-los em empregos",
      "Glass ceiling de progressao: gap cresce\nde 9,1% (jovens) para 37,5% (35-44 anos)"]),
]
for i, (color, tag, title, items) in enumerate(barreiras):
    x = In(0.3) + i * In(4.35)
    add_rect(s, x, In(1.62), In(4.1), In(5.15),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=2)
    add_rect(s, x, In(1.62), In(4.1), In(0.92), fill_rgb=color)
    add_text(s, tag, x+In(0.15), In(1.65), In(3.8), In(0.45),
             font_size=14, bold=True, color=C_WHITE)
    add_text(s, title, x+In(0.15), In(2.08), In(3.8), In(0.4),
             font_size=12.5, color=RGBColor(0xFF,0xF5,0xE0) if color==C_AMBER else RGBColor(0xBB,0xDE,0xFB),
             italic=True)
    for j, item in enumerate(items):
        add_text(s, "▸  " + item, x+In(0.2), In(2.65)+j*In(0.8), In(3.7), In(0.75),
                 font_size=12.5, color=C_BLACK)

add_rect(s, In(0.3), In(6.87), In(12.7), In(0.38),
         fill_rgb=RGBColor(0xFF,0xF9,0xE7), line_rgb=C_AMBER, line_pt=1)
add_text(s, "Conclusao: politicas que agem apenas no salario atacam 16% do problema. 84% esta na porta de entrada.",
         In(0.5), In(6.9), In(12.3), In(0.32),
         font_size=12.5, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
footer(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — O GAP: DO BRUTO AO LÍQUIDO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "3. O Gap Racial — Do Bruto ao Residual",
           "HLM 3 niveis (N=7,7M) decompoe o gap em camadas causais")

add_text(s, "Quanto do gap salarial e discriminacao pura?",
         In(0.4), In(1.12), In(12.5), In(0.42),
         font_size=18, bold=True, color=C_DARK)

levels = [
    ("Gap BRUTO",          19.3, C_RED,   "Sem controles — diferenca observada"),
    ("Mediacao UPA (bairro)", 9.6, C_BLUE,  "52% explicado por onde o trabalhador mora"),
    ("Gap LIQUIDO",        9.7,  C_AMBER, "Discriminacao residual pos-contexto"),
    ("Mediacao Ocupacional", 3.5, C_GREEN, "18% explicado pelo tipo de funcao"),
    ("Gap PURO",           6.2,  C_RED,   "Discriminacao de remuneracao irredutivel"),
]
bar_top  = In(1.68)
bar_left = In(3.8)
bar_w_total = In(9.0)
max_val  = 20.0

for i, (label, val, color, note) in enumerate(levels):
    bw = bar_w_total * (val / max_val)
    add_rect(s, bar_left, bar_top + i*In(0.88), bw, In(0.68), fill_rgb=color)
    add_text(s, label, In(0.2), bar_top + i*In(0.88) + In(0.12),
             In(3.5), In(0.5), font_size=13, bold=True, color=color, align=PP_ALIGN.RIGHT)
    add_text(s, f"{val:.1f}%", bar_left + bw + In(0.12),
             bar_top + i*In(0.88) + In(0.14), In(0.8), In(0.42),
             font_size=18, bold=True, color=color)
    add_text(s, note, bar_left + bw + In(1.05),
             bar_top + i*In(0.88) + In(0.14), In(4.5), In(0.42),
             font_size=11, color=C_GRAY, italic=True)

add_rect(s, In(0.3), In(6.35), In(12.7), In(0.65),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=C_RED, line_pt=1.5)
add_text(s,
    "70% do gap e explicado por contexto (bairro) e ocupacao — mas 6,2% persiste como discriminacao pura,\n"
    "mesmo controlando escolaridade, experiencia, setor, formalidade e funcao.",
    In(0.5), In(6.4), In(12.3), In(0.55),
    font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
footer(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BARREIRA DE ACESSO (GLMM + OB)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "4. Barreira de Acesso — Onde a Discriminacao Comeca",
           "Regressao logistica multinivel (N=7,7M, 40,9k bairros) + Oaxaca-Blinder")

# Lado esquerdo: números impactantes
add_text(s, "Barreira de ACESSO a empregos qualificados:",
         In(0.4), In(1.15), In(6.5), In(0.45),
         font_size=17, bold=True, color=C_DARK)

big_kpi(s, "Odds Ratio (GLMM M2)",
        or_str(P["OR_M2"]),
        "chance de ocupacao qualificada (M2 — apos todos os controles)",
        In(0.4), In(1.7), w=In(6.1), h=In(1.7),
        val_color=C_RED, bg=RGBColor(0x1A,0x05,0x05), border=C_RED)

big_kpi(s, "AME — efeito marginal medio (M2)",
        ame(P["AME_M2_pp"]),
        "gap de probabilidade de acesso que persiste com bairro identico",
        In(0.4), In(3.55), w=In(6.1), h=In(1.7),
        val_color=C_AMBER, bg=RGBColor(0x1A,0x15,0x00), border=C_AMBER)

add_rect(s, In(0.4), In(5.35), In(6.1), In(1.1),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=C_RED, line_pt=1)
add_text(s,
    f"Mesmo com EDUCACAO, SEXO, IDADE e BAIRRO\nIDENTICOS aos de um branco, um trabalhador\nnegro tem {fmt(P['OR_M2_menor_pct'],1)}% MENOS chance de trabalho qualificado.",
    In(0.55), In(5.4), In(5.8), In(1.0),
    font_size=13, bold=True, color=C_RED)

# Lado direito: Oaxaca-Blinder
add_text(s, "Onde opera o gap? Oaxaca-Blinder:",
         In(7.1), In(1.15), In(5.9), In(0.45),
         font_size=17, bold=True, color=C_DARK)

ob_items = [
    (In(7.1), In(1.7), In(5.9), In(2.3), C_BLUE,
     "84%", "Efeito ACESSO (dotacoes)",
     "Negros estao em ocupacoes piores,\nsetores menos remunerados,\nemprego informal. Discriminacao\npre-mercado e de recrutamento."),
    (In(7.1), In(4.1), In(5.9), In(1.9), C_RED,
     "16%", "Efeito SALARIO (retornos)",
     "O mercado remunera as mesmas\ncaracteristicas a taxas menores\npara negros — discriminacao\ndentro da funcao."),
]
for l, t, w, h, color, pct, title, body in ob_items:
    add_rect(s, l, t, w, h, fill_rgb=C_LGRAY, line_rgb=color, line_pt=1.5)
    add_rect(s, l, t, In(1.1), h, fill_rgb=color)
    add_text(s, pct, l+In(0.1), t+In(0.35), In(0.9), In(0.7),
             font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, l+In(1.2), t+In(0.1), w-In(1.35), In(0.42),
             font_size=14, bold=True, color=color)
    add_text(s, body, l+In(1.2), t+In(0.55), w-In(1.35), h-In(0.65),
             font_size=12, color=C_BLACK)

footer(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GLASS CEILING E CICLO DE VIDA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "5. Glass Ceiling — A Barreira Cresce com a Carreira",
           "Regressao quantilica + curva de ciclo de vida racial | gap piora ao longo do tempo")

add_img(s, FIGURES / "quantreg_trajetoria.png", In(0.3), In(1.15), In(6.8), In(4.3))

add_text(s, "O que os dados mostram:",
         In(7.5), In(1.15), In(5.6), In(0.42),
         font_size=17, bold=True, color=C_DARK)

insights = [
    (C_RED,   "Gap sobe no topo",
     "Gap de 8,0% na mediana (q50)\npara 11,8% no topo (q90)\n— glass ceiling formal confirmado\n(QR global sem controles ocupacionais)"),
    (C_AMBER, "Gap piora com a idade",
     "14-24 anos: gap = 9,1%\n35-44 anos: gap = 37,5%\nBarreira aprofunda-se na carreira"),
    (C_BLUE,  "Causa: dotacoes, nao retornos",
     "RIF-OB: q90 = 75,8% por ACESSO\ne 11,2% por discriminacao salarial\nFalha esta em antes do emprego"),
]
for i, (color, title, body) in enumerate(insights):
    add_rect(s, In(7.5), In(1.68)+i*In(1.52), In(5.6), In(1.38),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=1.5)
    add_rect(s, In(7.5), In(1.68)+i*In(1.52), In(5.6), In(0.42), fill_rgb=color)
    add_text(s, title, In(7.65), In(1.72)+i*In(1.52), In(5.3), In(0.38),
             font_size=13, bold=True, color=C_WHITE)
    add_text(s, body, In(7.65), In(2.15)+i*In(1.52), In(5.3), In(0.85),
             font_size=12.5, color=C_BLACK)

add_rect(s, In(0.3), In(5.55), In(12.7), In(0.7),
         fill_rgb=RGBColor(0xFF,0xF9,0xE7), line_rgb=C_AMBER, line_pt=1)
add_text(s,
    "Sticky Floor: discriminacao e 3x mais intensa na base (q10=33,1%) do que no topo (q90=11,2%).\n"
    "Enforcement anti-discriminatorio deve focar nos trabalhadores de menor renda.",
    In(0.5), In(5.6), In(12.3), In(0.6),
    font_size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
footer(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — REDES E DISCRIMINACAO GEOGRAFICA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "6. Redes e Heterogeneidade Geografica — A Dimensao Invisivel",
           "SNA: exclusao das redes | HLM Random Slope: discriminacao varia 5x entre estados (N=7,7M)")

# Painel SNA (esquerdo)
add_text(s, "Exclusao das Redes (SNA):",
         In(0.3), In(1.12), In(6.2), In(0.42),
         font_size=16, bold=True, color=C_DARK)

add_img(s, FIGURES / "sna_rede_demografica.png", In(0.3), In(1.6), In(6.2), In(3.8))

add_rect(s, In(0.3), In(5.5), In(6.2), In(1.1),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=C_RED, line_pt=1)
add_text(s,
    "Betweenness centrality = 0 para negros\nem TODOS os niveis de escolaridade.\nDiplomas negros valem menos porque\nfaltam redes para converter titulos em empregos.",
    In(0.45), In(5.55), In(5.9), In(1.0),
    font_size=13, bold=True, color=C_RED)

# Painel Random Slope (direito)
add_text(s, "Discriminacao varia entre estados (HLM):",
         In(6.9), In(1.12), In(6.1), In(0.42),
         font_size=16, bold=True, color=C_DARK)

rs_b    = P.get("RS_B_NEGRO_FIXO", -0.1218)
rs_gap  = abs(P.get("RS_GAP_PCT", -11.5))
rs_lo   = P.get("RS_GAP_LO_PCT",  -20.9)
rs_hi   = P.get("RS_GAP_HI_PCT",  -2.0)
rs_sd   = P.get("RS_SD_NEGRO", 0.0506)
rs_rho  = P.get("RS_RHO", -0.3719)
rs_lr   = P.get("RS_LRT_LR", 7611.6)

big_kpi(s, "Gap medio (efeito fixo)",
        f"{rs_gap:.1f}%",
        "penalidade salarial racial media nacional",
        In(6.9), In(1.65), w=In(5.9), h=In(1.4),
        val_color=C_RED, bg=RGBColor(0x1A,0x05,0x05), border=C_RED)

add_rect(s, In(6.9), In(3.18), In(5.9), In(1.15),
         fill_rgb=C_LGRAY, line_rgb=C_AMBER, line_pt=1.2)
add_text(s, "Variacao entre estados (95% das UFs):",
         In(7.05), In(3.22), In(5.6), In(0.35),
         font_size=12, bold=True, color=C_AMBER)
add_text(s, f"de {rs_lo:.1f}% ate {rs_hi:.1f}% — amplitude de {abs(rs_lo-rs_hi):.1f} pp",
         In(7.05), In(3.6), In(5.6), In(0.55),
         font_size=22, bold=True, color=C_DARK)

add_rect(s, In(6.9), In(4.45), In(5.9), In(1.15),
         fill_rgb=C_LGRAY, line_rgb=C_BLUE, line_pt=1.2)
add_text(s, f"Correlacao estado rico x gap racial: rho = {rs_rho:.3f}",
         In(7.05), In(4.49), In(5.6), In(0.35),
         font_size=12, bold=True, color=C_BLUE)
add_text(s,
    "Estados mais pobres concentram maior penalidade racial.\n"
    "Politicas nacionais uniformes ignoram essa heterogeneidade.",
    In(7.05), In(4.87), In(5.6), In(0.68),
    font_size=12.5, color=C_BLACK)

add_rect(s, In(6.9), In(5.72), In(5.9), In(0.88),
         fill_rgb=RGBColor(0xE3,0xF2,0xFD), line_rgb=C_DARK, line_pt=1)
add_text(s, f"LRT boundary test: LR={rs_lr:.0f} | p<0,001*** — variacao geografica CONFIRMADA (N=7.689.426)",
         In(7.05), In(5.77), In(5.6), In(0.75),
         font_size=12, bold=True, color=C_DARK)

footer(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ROBUSTEZ: O GAP E REAL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "7. O Gap e Real — Tres Testes de Robustez",
           "Nenhuma variavel omitida plausivel consegue anular os resultados")

add_text(s, "Poderia o gap ser explicado por variaveis nao observadas?",
         In(0.4), In(1.12), In(12.5), In(0.42),
         font_size=18, bold=True, color=C_DARK)

rob_items = [
    (C_DARK, "Konfound (Frank, 2013)",
     "99,5% dos casos",
     "precisariam ser removidos ou alterados para anular o resultado",
     "Como referencia: um confundidor assim nunca existiu em Ciencias Sociais"),
    (C_RED,  "E-values (VanderWeele, 2017)",
     f"E >= {fmt(P['EVAL_M2'],2)}x",
     "de associacao com raca E com ocupacao simultaneamente",
     "Uma variavel oculta teria que ser 2,3x mais preditiva que TODAS as observadas juntas"),
    (C_BLUE, "Oster Bounds (Oster, 2019)",
     "delta* negativo",
     "variaveis omitidas teriam que agir na direcao OPOSTA as observadas",
     "Impossivel: estudos mostram que omitidas e observadas agem na mesma direcao"),
]
for i, (color, metodo, val, interpretacao, nota) in enumerate(rob_items):
    add_rect(s, In(0.3)+i*In(4.35), In(1.65), In(4.1), In(4.85),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=2)
    add_rect(s, In(0.3)+i*In(4.35), In(1.65), In(4.1), In(0.5), fill_rgb=color)
    add_text(s, metodo, In(0.45)+i*In(4.35), In(1.69), In(3.8), In(0.45),
             font_size=13, bold=True, color=C_WHITE)
    add_text(s, val, In(0.45)+i*In(4.35), In(2.28), In(3.8), In(0.65),
             font_size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, interpretacao, In(0.45)+i*In(4.35), In(3.0), In(3.8), In(1.1),
             font_size=13, color=C_BLACK)
    add_rect(s, In(0.35)+i*In(4.35), In(4.18), In(4.0), In(0.03), fill_rgb=C_GRAY)
    add_text(s, nota, In(0.45)+i*In(4.35), In(4.25), In(3.8), In(1.1),
             font_size=11.5, italic=True, color=C_GRAY)

add_rect(s, In(0.3), In(6.62), In(12.7), In(0.58),
         fill_rgb=RGBColor(0x1F,0x38,0x64))
add_text(s,
    "Veredicto: os tres metodos independentes concordam. O gap racial no mercado de trabalho brasileiro e real,\nrobusto e nao pode ser explicado por variaveis ocultas plausíveis.",
    In(0.5), In(6.67), In(12.3), In(0.48),
    font_size=13.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
footer(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — O QUE FAZER: POLITICAS PRIORIZADAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "8. O Que Fazer — Tres Eixos de Politica Publica",
           "TOPSIS + AHP + Programacao Linear — prioridade: atacar a porta de entrada, nao o salario")

add_text(s, "84% do gap e de ACESSO. As politicas precisam refletir isso.",
         In(0.4), In(1.1), In(12.5), In(0.42),
         font_size=17, bold=True, color=C_DARK)

eixos = [
    (C_RED,   "#1 — Acesso",
     [("Cotas CBO (Dirigente/Profissional)", "meta IR >= 0,80 ate 2030",
       f"TOPSIS CC={fmt(P['TOPSIS_P1_CC'],3)} — 2,1x superior ao 2o colocado"),
      ("PRONATEC c/ recorte racial", "subsidio de transporte em UPAs pobres",
       "ataca segregacao residencial (52% do gap)"),
      ("Residencias em empresas de alta renda", "bolsas para egressos negros",
       "quebra a barreira de redes (betweenness=0)")]),
    (C_BLUE,  "#2 — Remuneracao",
     [("Transparencia salarial por raca/genero", "obrigatoria para empresas >100 func.",
       "visa o gap residual de 6,2% (HLM M4)"),
      ("Auditoria de igual pagamento", "penalidades progressivas",
       "enforcement direto do sticky floor (q10=33%)"),
      ("Piso salarial nas categorias de maior gap", "indexado ao gap HLM M4",
       "protecao prioritaria para trabalhadores da base")]),
    (C_DARK,  "#3 — Redes",
     [("Mentoria estruturada por empresas", "elevacao da betweenness centrality",
       "SNA: betweenness=0 anula retorno ao diploma"),
      ("30% DAS e liderancas corporativas", "reserva para negros ate 2030",
       "inclusao nas redes de decisao (glass ceiling)"),
      ("Equidade educacional de qualidade", "nao apenas de acesso",
       f"gap qualificacao: 11% negros vs 22% brancos — estavel em 10 anos")]),
]
for i, (color, titulo, acoes) in enumerate(eixos):
    x = In(0.3) + i * In(4.35)
    add_rect(s, x, In(1.62), In(4.1), In(5.55),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=2)
    add_rect(s, x, In(1.62), In(4.1), In(0.5), fill_rgb=color)
    add_text(s, titulo, x+In(0.15), In(1.66), In(3.8), In(0.45),
             font_size=14, bold=True, color=C_WHITE)
    for j, (acao, sub, evidencia) in enumerate(acoes):
        y = In(2.22) + j * In(1.6)
        add_rect(s, x+In(0.15), y, In(3.8), In(1.45),
                 fill_rgb=C_WHITE, line_rgb=color, line_pt=0.5)
        add_text(s, acao, x+In(0.25), y+In(0.06), In(3.6), In(0.42),
                 font_size=12.5, bold=True, color=color)
        add_text(s, sub, x+In(0.25), y+In(0.48), In(3.6), In(0.38),
                 font_size=11.5, color=C_BLACK)
        add_text(s, evidencia, x+In(0.25), y+In(0.88), In(3.6), In(0.42),
                 font_size=10.5, italic=True, color=C_GRAY)

add_rect(s, In(0.3), H-In(0.6), In(12.7), In(0.32), fill_rgb=C_AMBER)
add_text(s, "Pareto (lambda=0,5): combinacao otima = Eixo 1 (cotas) + Eixo 3 (mentoria). Atuacao isolada em remuneracao e insuficiente.",
         In(0.5), H-In(0.57), In(12.3), In(0.28),
         font_size=11.5, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
footer(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSAO EXECUTIVA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=C_DARK)
add_rect(s, 0, 0, W, In(0.75), fill_rgb=RGBColor(0x0D,0x1F,0x3C))
add_text(s, "9. Conclusao — O Que Este Trabalho Prova",
         In(0.4), In(0.1), In(12.5), In(0.6),
         font_size=25, bold=True, color=C_WHITE)

# KPIs superiores
numeros = [
    (f"{fmt(P.get('GB_HLM_M1_pct', 19.3),1)}%", "Gap bruto",       "sem controles (HLM M1)"),
    ("6,2%",          "Discriminacao pura",  "mesma funcao e empresa"),
    (or_str(P["OR_M2"]), "Odds Ratio GLMM", "acesso ocupacao qualificada"),
    ("0",             "Betweenness negros",  "em todas as escolaridades (SNA)"),
]
for i, (val, label, sub) in enumerate(numeros):
    x = In(0.4) + i * In(3.2)
    add_rect(s, x, In(0.85), In(3.0), In(1.4),
             fill_rgb=RGBColor(0x1A,0x30,0x5C), line_rgb=C_AMBER, line_pt=1)
    add_text(s, val, x+In(0.1), In(0.9), In(2.8), In(0.7),
             font_size=28, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_text(s, label, x+In(0.1), In(1.6), In(2.8), In(0.35),
             font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, sub, x+In(0.1), In(1.95), In(2.8), In(0.25),
             font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

# Mensagens-chave
mensagens = [
    (C_RED,   "A discriminacao e real e mensuravel",
     "5 metodos independentes chegam ao mesmo diagnostico. Nao e artefato estatistico."),
    (C_AMBER, "Opera em duas etapas distintas",
     "84% na PORTA DE ENTRADA (acesso a funcoes qualificadas) + 16% no SALARIO dentro da funcao."),
    (C_BLUE,  "E geograficamente heterogenea",
     f"O gap varia de {P.get('RS_GAP_LO_PCT',-19.8):.1f}% a {P.get('RS_GAP_HI_PCT',-2.2):.1f}% entre estados (random slope HLM, N={fmtN(P.get('RS_N_OBS',7689426))}). Estados pobres concentram mais discriminacao."),
    (C_GREEN, "Politicas precisam ser especificas",
     "Cotas ocupacionais (CBO) + mentoria de redes atacam 84% do problema. Transparencia salarial ataca o restante."),
]
for i, (color, titulo, corpo) in enumerate(mensagens):
    add_rect(s, In(0.3), In(2.42)+i*In(1.02), In(12.7), In(0.94),
             fill_rgb=RGBColor(0x0F,0x27,0x4A),
             line_rgb=color, line_pt=1.5)
    add_rect(s, In(0.3), In(2.42)+i*In(1.02), In(0.18), In(0.94), fill_rgb=color)
    add_text(s, titulo, In(0.6), In(2.46)+i*In(1.02), In(4.5), In(0.4),
             font_size=14, bold=True, color=color)
    add_text(s, corpo, In(0.6), In(2.86)+i*In(1.02), In(12.2), In(0.45),
             font_size=13, color=C_WHITE)

add_rect(s, 0, H-In(0.45), W, In(0.45), fill_rgb=RGBColor(0x0D,0x1F,0x3C))
add_text(s,
    "Ricardo Calheiros  |  MBA Data Science & Analytics  |  USP/ESALQ  |  rickinrj@gmail.com",
    In(0.3), H-In(0.42), In(12.7), In(0.38),
    font_size=11, color=RGBColor(0x90,0xA4,0xAE), align=PP_ALIGN.CENTER)

# ── Salvar ────────────────────────────────────────────────────────────────────
prs.save(str(OUT_PPT))
print(f"Arquivo gerado: {OUT_PPT.name}")
print(f"Slides: {len(prs.slides)}")
print(f"Tamanho: {OUT_PPT.stat().st_size // 1024} KB")
print("Abra com PowerPoint para revisar.")
