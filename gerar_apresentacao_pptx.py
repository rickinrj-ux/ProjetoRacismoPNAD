"""
gerar_apresentacao_pptx.py
Gera apresentacao_tcc.pptx — 15 slides para defesa de TCC na banca.
Usa python-pptx (pip install python-pptx).
"""
import sys; sys.stdout.reconfigure(encoding='utf-8')
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as In, Pt
import pandas as pd
import numpy as np

ROOT    = Path(r"C:\Users\user\Documents\ProjetoRacismoPNAD")
FIGURES = ROOT / "outputs" / "figures"
TABLES  = ROOT / "outputs" / "tables"
OUT_PPT = ROOT / "apresentacao_tcc.pptx"

# ── Paleta ────────────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1F, 0x38, 0x64)   # azul escuro (cabeçalhos)
C_BLUE   = RGBColor(0x15, 0x65, 0xC0)   # azul (brancos / positivo)
C_RED    = RGBColor(0xB7, 0x1C, 0x1C)   # vermelho (negros / discriminação)
C_AMBER  = RGBColor(0xFF, 0x8F, 0x00)   # âmbar (destaque)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x21, 0x21, 0x21)
C_GRAY   = RGBColor(0x61, 0x61, 0x61)
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)

W = In(13.33)   # largura slide 16:9
H = In(7.50)    # altura slide 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # layout em branco

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = _Pt(line_pt) if line_pt else 0
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=20, bold=False, italic=False,
             color=C_BLACK, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tf

def add_multiline(slide, lines, l, t, w, h,
                  font_size=16, color=C_BLACK, bold=False,
                  align=PP_ALIGN.LEFT, line_spacing=1.15,
                  font_name="Calibri"):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = True
    first = True
    for line_text in lines:
        if first:
            p = tf.text_frame.paragraphs[0]
            first = False
        else:
            p = tf.text_frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tf

def add_img(slide, img_path, l, t, w, h=None):
    p = Path(img_path)
    if not p.exists():
        add_text(slide, f"[Figura: {p.name}]", l, t, w, In(1),
                 font_size=10, color=C_GRAY, italic=True)
        return
    if h:
        slide.shapes.add_picture(str(p), l, t, w, h)
    else:
        slide.shapes.add_picture(str(p), l, t, w)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, In(1.10), fill_rgb=C_DARK)
    add_text(slide, title, In(0.3), In(0.08), In(12), In(0.55),
             font_size=26, bold=True, color=C_WHITE, font_name="Calibri")
    if subtitle:
        add_text(slide, subtitle, In(0.3), In(0.65), In(12), In(0.38),
                 font_size=14, color=RGBColor(0xBB, 0xDE, 0xFB), font_name="Calibri")

def bullet_box(slide, items, l, t, w, h,
               dot="▸", font_size=16, color=C_BLACK,
               dot_color=C_BLUE):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.text_frame.paragraphs[0]
            first = False
        else:
            p = tf.text_frame.add_paragraph()
        p.space_after = Pt(4)
        # dot
        r1 = p.add_run(); r1.text = dot + "  "
        r1.font.size = Pt(font_size); r1.font.color.rgb = dot_color
        r1.font.bold = True; r1.font.name = "Calibri"
        # text
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(font_size); r2.font.color.rgb = color
        r2.font.name = "Calibri"
    return tf

def kpi_box(slide, label, value, unit, l, t, w=In(2.8), h=In(1.3),
            val_color=C_BLUE):
    add_rect(slide, l, t, w, h, fill_rgb=C_LGRAY,
             line_rgb=C_DARK, line_pt=1)
    add_text(slide, label, l+In(0.1), t+In(0.05), w-In(0.2), In(0.35),
             font_size=11, color=C_GRAY, font_name="Calibri")
    add_text(slide, value, l+In(0.1), t+In(0.35), w-In(0.2), In(0.6),
             font_size=30, bold=True, color=val_color, font_name="Calibri")
    add_text(slide, unit, l+In(0.1), t+In(0.92), w-In(0.2), In(0.35),
             font_size=10, color=C_GRAY, italic=True, font_name="Calibri")

def footer(slide, slide_num, total=25):
    add_rect(slide, 0, H-In(0.28), W, In(0.28), fill_rgb=C_DARK)
    add_text(slide,
             "Ricardo Calheiros  |  MBA USP/ESALQ  |  Racismo Estrutural e Mercado de Trabalho",
             In(0.15), H-In(0.25), In(11), In(0.25),
             font_size=8.5, color=RGBColor(0xBB,0xDE,0xFB), font_name="Calibri")
    add_text(slide, f"{slide_num}/{total}",
             In(12.8), H-In(0.25), In(0.5), In(0.25),
             font_size=8.5, color=C_WHITE, align=PP_ALIGN.RIGHT, font_name="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=C_DARK)
add_rect(s, In(0.5), In(0.4), In(12.33), In(3.8),
         fill_rgb=RGBColor(0x0D,0x1F,0x3C), line_rgb=C_AMBER, line_pt=1.5)

add_text(s, "Racismo Estrutural e Mercado de Trabalho no Brasil",
         In(0.9), In(0.65), In(11.5), In(1.5),
         font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(s, "Evidências da PNAD Contínua 2016–2025 via HLM, ML/SHAP,\nOaxaca-Blinder, Logit Multinível e Regressão Quantílica",
         In(0.9), In(2.15), In(11.5), In(0.9),
         font_size=16, color=RGBColor(0xBB,0xDE,0xFB), align=PP_ALIGN.CENTER, font_name="Calibri")
add_rect(s, In(0.9), In(3.1), In(11.5), In(0.03), fill_rgb=C_AMBER)

add_text(s, "Ricardo Calheiros",
         In(0.9), In(3.3), In(11.5), In(0.45),
         font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "MBA em Data Science e Analytics  |  USP/ESALQ  |  Defesa: 2026",
         In(0.9), In(3.75), In(11.5), In(0.4),
         font_size=13, color=RGBColor(0xBB,0xDE,0xFB), align=PP_ALIGN.CENTER)

# stats rápidas no rodapé da capa
for i, (val, lbl) in enumerate([
    ("15,9M", "observações brutas"),
    ("7,7M",  "obs. PEA c/ renda"),
    ("5",     "metodologias"),
    ("10",    "anos de série"),
]):
    x = In(0.5) + i * In(3.2)
    add_rect(s, x, In(4.7), In(2.9), In(1.0),
             fill_rgb=RGBColor(0x1A,0x30,0x5C), line_rgb=C_AMBER, line_pt=0.8)
    add_text(s, val, x+In(0.1), In(4.72), In(2.7), In(0.55),
             font_size=28, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x+In(0.1), In(5.28), In(2.7), In(0.35),
             font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Orientador(a): _______________",
         In(0.3), H-In(0.6), In(8), In(0.4),
         font_size=12, color=RGBColor(0x90,0xA4,0xAE), italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O PROBLEMA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "1. O Problema",
           "Por que um mestrado quantitativo sobre o gap racial brasileiro?")

add_text(s, "O gap salarial racial persiste mesmo após décadas de políticas inclusivas",
         In(0.4), In(1.25), In(12.5), In(0.5),
         font_size=18, bold=True, color=C_DARK)

bullet_box(s, [
    "Trabalhadores negros ganham em média 40–55% menos que brancos — bruta, sem controles",
    "O gap persiste mesmo após controlar escolaridade, sexo e experiência: é estrutural",
    "Mecanismos: segregação residencial, exclusão das redes profissionais, glass ceiling ocupacional",
    "Estudos tradicionais usam apenas OLS — não separam ACESSO de REMUNERAÇÃO",
], In(0.4), In(1.85), In(7.2), In(2.8), font_size=15)

add_rect(s, In(7.8), In(1.2), In(5.2), In(4.8),
         fill_rgb=C_LGRAY, line_rgb=C_DARK, line_pt=0.8)
add_text(s, "Inovações deste TCC", In(8.0), In(1.35), In(4.8), In(0.4),
         font_size=14, bold=True, color=C_DARK)
for i, item in enumerate([
    "✦ 15,9M obs. — série histórica completa",
    "✦ Variáveis CBO, vínculo e horas (VD4009/V4010/VD4031) extraídas dos ZIPs brutos do IBGE",
    "✦ Discriminação de ACESSO (logit)",
    "✦ Discriminação de REMUNERAÇÃO (HLM + QR)",
    "✦ Convergência de 5 métodos independentes",
]):
    add_text(s, item, In(8.1), In(1.85)+i*In(0.65), In(4.7), In(0.6),
             font_size=13, color=C_BLACK if i > 0 else C_BLUE)

footer(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DADOS E DATASET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "2. Base de Dados", "PNAD Contínua 2016–2025 + enriquecimento com variáveis ocupacionais")

kpi_box(s, "Observações brutas", "15,9M", "PNAD 2016–2025", In(0.3), In(1.2))
kpi_box(s, "PEA com renda positiva", "7,7M", "após filtros", In(3.3), In(1.2))
kpi_box(s, "UPAs (bairros)", "40.120", "Nível 2 do HLM", In(6.3), In(1.2))
kpi_box(s, "Cobertura temporal", "10 anos", "2016 T1 → 2025", In(9.3), In(1.2))

add_text(s, "Enriquecimento dos microdados:", In(0.4), In(2.7), In(12), In(0.4),
         font_size=16, bold=True, color=C_DARK)

cols = [
    ("V4010  →  ocp_grupo_cbo", "Grupo ocupacional CBO-Domiciliar\n(10 grupos ISCO-08)"),
    ("VD4009  →  emprego_formal", "Vínculo empregatício detalhado\n(carteira, servidor, conta-própria...)"),
    ("VD4031  →  horas_c", "Total de horas efetivamente\ntrabalhadas (centralizada)"),
]
for i, (var, desc) in enumerate(cols):
    x = In(0.3) + i*In(4.3)
    add_rect(s, x, In(3.15), In(4.1), In(1.5),
             fill_rgb=RGBColor(0xE3,0xF2,0xFD), line_rgb=C_BLUE, line_pt=0.8)
    add_text(s, var, x+In(0.15), In(3.22), In(3.8), In(0.5),
             font_size=13, bold=True, color=C_BLUE, font_name="Courier New")
    add_text(s, desc, x+In(0.15), In(3.7), In(3.8), In(0.85),
             font_size=12, color=C_BLACK)

add_text(s, "Estratégia: leitura seletiva dos ZIPs originais (≈8,4 GB) — sem re-download.\n"
            "Join por chave composta (Ano+Trimestre+UPA+V1008+V2003) → taxa NA = 0% na amostra de trabalho.",
         In(0.4), In(4.85), In(12.5), In(0.9),
         font_size=13, color=C_GRAY, italic=True)

footer(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ANÁLISE DESCRITIVA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "3. Análise Descritiva — A Desigualdade nos Dados",
           "Gap mediano bruto de 27,5% | 7,7 milhões de empregados com renda positiva | PNAD Contínua 2016–2025")

add_img(s, FIGURES / "fig1_densidade_log_salario.png", In(0.3), In(1.2), In(6.4), In(4.5))

add_text(s, "O que os números dizem", In(7.0), In(1.2), In(6.1), In(0.4),
         font_size=16, bold=True, color=C_DARK)

kpi_box(s, "Gap mediano bruto",     "27,5%",    "brancos R$2.000 vs negros R$1.450",  In(7.0),  In(1.75), w=In(2.9), val_color=C_RED)
kpi_box(s, "Gap na média",          "37,5%",    "brancos R$3.256 vs negros R$2.040",  In(10.1), In(1.75), w=In(3.1), val_color=C_RED)
kpi_box(s, "Gap log-renda",         "42,0%",    "equivalente em renda esperada",      In(7.0),  In(3.15), w=In(2.9), val_color=C_RED)
kpi_box(s, "Formal (brancos/neg.)", "53%/48%",  "5 p.p. a menos de emprego formal",  In(10.1), In(3.15), w=In(3.1), val_color=C_AMBER)

add_rect(s, In(7.0), In(4.6), In(6.1), In(1.55),
         fill_rgb=RGBColor(0xE3,0xF2,0xFD), line_rgb=C_DARK, line_pt=0.8)
add_text(s, "Interpretação da curva de densidade:",
         In(7.15), In(4.65), In(5.8), In(0.35),
         font_size=12, bold=True, color=C_DARK)
add_multiline(s, [
    "▸  Curva vermelha (negros) deslocada à esquerda ao longo de toda a distribuição",
    "▸  Linhas tracejadas = médias de log-renda de cada grupo",
    "▸  O deslocamento é estrutural — não se concentra em nenhuma faixa específica",
], In(7.15), In(5.05), In(5.8), In(1.0), font_size=11, color=C_BLACK)

add_rect(s, In(0.3), In(5.85), In(12.7), In(0.45),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=C_RED, line_pt=0.8)
add_text(s, "O gap persiste em TODOS os níveis educacionais — inclusive pós-graduação (19,1%). "
            "Educação é condição necessária, mas não suficiente para eliminar a desigualdade racial.",
         In(0.5), In(5.9), In(12.3), In(0.4),
         font_size=12, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
footer(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARQUITETURA METODOLÓGICA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "4. Arquitetura Metodológica",
           "5 metodologias complementares — convergência de evidências")

metodos = [
    (C_DARK,  "HLM",            "Modelos Hierárquicos\nLineares (3 níveis)",        "Gap médio e\nmediação contextual"),
    (C_BLUE,  "ML/SHAP",        "XGBoost + Random Forest\n+ SHAP Values",           "Importância e\ndireção das variáveis"),
    (C_RED,   "Oaxaca-Blinder", "Decomposição do gap\nem dotações vs retornos",     "ONDE opera\na discriminação"),
    (C_AMBER, "Logit Multinível","Regressão logística\ncom estrutura multinível",   "Gap de ACESSO\na oportunidades"),
    (C_GREEN, "Reg. Quantílica", "Gap por quantil\nda distribuição",                "Glass Ceiling\nformal"),
]
for i, (color, title, desc, purpose) in enumerate(metodos):
    x = In(0.25) + i * In(2.6)
    add_rect(s, x, In(1.25), In(2.45), In(4.8),
             fill_rgb=RGBColor(0xF5,0xF5,0xF5), line_rgb=color, line_pt=1.5)
    add_rect(s, x, In(1.25), In(2.45), In(0.5), fill_rgb=color)
    add_text(s, title, x+In(0.1), In(1.28), In(2.25), In(0.45),
             font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+In(0.1), In(1.85), In(2.25), In(0.9),
             font_size=11.5, color=C_BLACK, align=PP_ALIGN.CENTER)
    add_rect(s, x+In(0.1), In(2.85), In(2.25), In(0.02), fill_rgb=color)
    add_text(s, purpose, x+In(0.1), In(2.95), In(2.25), In(0.7),
             font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

add_text(s, "↓ Convergência dos 5 métodos = robustez diagnóstica ↓",
         In(0.3), In(6.1), In(12.7), In(0.45),
         font_size=15, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
add_rect(s, In(0.3), In(6.15), In(12.7), In(0.5),
         fill_rgb=RGBColor(0xE8,0xEA,0xF0), line_rgb=C_DARK, line_pt=0.5)
add_text(s, "Discriminação de ACESSO (logit) + Discriminação de REMUNERAÇÃO (HLM + QR) + Exclusão de REDES (SNA)",
         In(0.4), In(6.2), In(12.5), In(0.4),
         font_size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
footer(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HLM: DECOMPOSIÇÃO DO GAP
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "5. HLM — Decomposição do Gap em Camadas",
           "Do gap bruto de 19,3% ao resíduo de discriminação pura de 6,2%")

# Barra de decomposição visual
levels = [
    ("Gap BRUTO\n(M1)", 19.3, C_RED,   "Sem controles contextuais\nnem ocupacionais"),
    ("Mediação UPA\n(M2)", 9.6, C_BLUE,  "Local de moradia explica\n52% do gap bruto"),
    ("Gap LÍQUIDO\n(M3)", 9.7, C_AMBER, "Após contexto UPA + UF FE\n(discriminação residual)"),
    ("Mediação Ocp.\n(M4)", 3.5, C_GREEN, "CBO + formalidade + horas\nexplicam mais 18,1%"),
    ("Gap PURO\n(M4)", 6.2, C_RED,   "Discriminação residual\npós-ocupação"),
]
bar_top = In(1.3)
bar_left = In(0.4)
bar_w_total = In(11.8)
max_val = 20.0

for i, (label, val, color, note) in enumerate(levels):
    bw = bar_w_total * (val / max_val)
    add_rect(s, bar_left, bar_top + i*In(0.97), bw, In(0.72), fill_rgb=color)
    add_text(s, f"{val:.1f}%", bar_left + bw + In(0.1), bar_top + i*In(0.97) + In(0.18),
             In(0.8), In(0.4), font_size=15, bold=True, color=color)
    add_text(s, label, In(0.05), bar_top + i*In(0.97) + In(0.05),
             In(1.4), In(0.65), font_size=10, color=C_BLACK, align=PP_ALIGN.RIGHT, bold=True)
    add_text(s, note, bar_left + bw + In(1.0), bar_top + i*In(0.97) + In(0.1),
             In(5.5), In(0.55), font_size=10, color=C_GRAY, italic=True)

add_rect(s, In(0.4), In(6.35), In(12.0), In(0.5),
         fill_rgb=RGBColor(0xFF,0xF9,0xE7), line_rgb=C_AMBER, line_pt=1)
add_text(s, "Achado central: 70,0% do gap é explicado por mediação contextual (UPA) e ocupacional — mas 6,2% persiste como discriminação de remuneração pura.",
         In(0.6), In(6.4), In(11.7), In(0.45),
         font_size=12.5, bold=True, color=C_DARK)
footer(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — COMPOSIÇÃO OCUPACIONAL E GLASS CEILING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "6. Composição Ocupacional — Onde Estão Negros e Brancos?",
           "Sub-representação sistemática nos grupos de alto prestígio")

add_img(s, FIGURES / "comp_razao_grupo_cbo.png", In(0.3), In(1.2), In(6.5))
add_img(s, FIGURES / "comp_representacao_topo.png", In(7.0), In(1.2), In(6.0))

add_rect(s, In(0.3), In(5.8), In(12.7), In(0.6),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=C_RED, line_pt=0.8)
add_text(s,
    "Dirigentes: apenas 42 negros por 100 brancos na mesma função  |  "
    "Top 5% das capitais: IR = 0,47 (negros são 47% do esperado)  |  "
    "Elementares: 2,12× mais negros que brancos",
    In(0.5), In(5.85), In(12.3), In(0.55),
    font_size=12, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
footer(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — OAXACA-BLINDER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "7. Oaxaca-Blinder — Onde Opera a Discriminação?",
           "Gap total 53,0% | 84% explicado por diferenças nas características dos trabalhadores")

add_img(s, FIGURES / "oaxaca_decomposicao.png", In(0.3), In(1.2), In(7.5))

add_text(s, "Interpretação", In(8.1), In(1.2), In(5.0), In(0.4),
         font_size=16, bold=True, color=C_DARK)
bullet_box(s, [
    "84% = Efeito DOTAÇÕES\nNegros têm menor acesso a ocupações de prestígio, emprego formal, mais horas em subemprego",
    "16% = Efeito RETORNOS\nO mercado remunera as mesmas características a taxas diferentes por raça",
    "Conclusão: discriminação opera primariamente no ACESSO, não no salário dentro da função",
], In(8.1), In(1.7), In(5.0), In(3.5), font_size=13, dot_color=C_BLUE)

add_rect(s, In(8.1), In(5.4), In(5.0), In(1.0),
         fill_rgb=RGBColor(0xE3,0xF2,0xFD), line_rgb=C_BLUE, line_pt=0.8)
add_text(s, "Implicação de política:", In(8.25), In(5.45), In(4.7), In(0.3),
         font_size=11, bold=True, color=C_BLUE)
add_text(s, "Combater APENAS desigualdade salarial é insuficiente. É preciso atacar as barreiras de ACESSO a ocupações qualificadas.",
         In(8.25), In(5.77), In(4.7), In(0.6),
         font_size=12, color=C_BLACK)
footer(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GAP POR SUBGRUPO, MINCER E CICLO DE VIDA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "8. Gap por Subgrupo e Ciclo de Vida",
           "Regressão Mincer com controles progressivos + curva de ciclo de vida racial")

add_img(s, FIGURES / "fig4_gap_faixa_etaria.png", In(0.3), In(1.2), In(6.2), In(4.5))

add_text(s, "Regressões Mincer progressivas", In(6.8), In(1.2), In(6.3), In(0.4),
         font_size=15, bold=True, color=C_DARK)

mincer_rows = [
    ("M1 — Só raça",                   "−23,1%", C_RED),
    ("M2 — + Sexo, educ., experiência", "−22,3%", C_RED),
    ("M3 — + Efeitos fixos de UF",      "−7,9%",  C_AMBER),
    ("M4 — + Horas, formal, setor",     "−7,0%",  C_AMBER),
    ("M5 — + Ocupação CBO",             "−6,2%",  C_GREEN),
]
for i, (spec, gap, color) in enumerate(mincer_rows):
    add_rect(s, In(6.8), In(1.72) + i * In(0.82), In(6.3), In(0.75),
             fill_rgb=RGBColor(0xFF,0xEB,0xEE) if color==C_RED else
                      (RGBColor(0xFF,0xF9,0xE7) if color==C_AMBER else RGBColor(0xE8,0xF5,0xE9)),
             line_rgb=color, line_pt=0.8)
    add_text(s, spec, In(6.95), In(1.76) + i * In(0.82), In(4.5), In(0.38),
             font_size=11, color=C_BLACK)
    add_text(s, gap,  In(11.5), In(1.76) + i * In(0.82), In(1.4), In(0.38),
             font_size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)

add_rect(s, In(0.3), In(5.85), In(12.7), In(0.9),
         fill_rgb=RGBColor(0xE3,0xF2,0xFD), line_rgb=C_DARK, line_pt=0.8)
add_text(s,
    "Ciclo de vida: gap de 9,1% (14–24 anos) cresce para 37,5% (35–44 anos) — a barreira racial "
    "se aprofunda ao longo da carreira.  |  Interação negro×experiência (β=−0,0019, p<0,001): "
    "o retorno à senioridade é menor para negros — glass ceiling de progressão longitudinal confirmado.",
    In(0.5), In(5.9), In(12.3), In(0.8),
    font_size=11.5, bold=True, color=C_DARK)
footer(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LOGIT MULTINÍVEL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "9. Logit Multinível — Gap de Oportunidades Racial",
           "GLMM lme4 com efeito aleatório de UPA — discriminação de ACESSO confirmada (n=2,4M, 39,8k UPAs)")

add_img(s, FIGURES / "glmm_odds_ratios_full.png", In(0.3), In(1.2), In(7.0))

add_text(s, "Resultados-chave (GLMM)", In(7.6), In(1.2), In(5.5), In(0.4),
         font_size=16, bold=True, color=C_DARK)

kpis_logit = [
    ("OR = 0,704", "Ocupação qualificada\n(M1 — só ind. + UPA RE)", C_RED),
    ("OR = 0,747", "Ocupação qualificada\n(M2 — +contexto UPA)", C_AMBER),
    ("AME = −1,30 p.p.", "Gap pós-controles individuais\n(M1 — sem contexto UPA)", C_RED),
    ("AME = −1,07 p.p.", "Gap residual c/ contexto UPA\n(M2 — discriminação líquida)", C_RED),
]
for i, (val, lbl, color) in enumerate(kpis_logit):
    add_rect(s, In(7.6), In(1.7)+i*In(1.1), In(5.4), In(0.9),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=1)
    add_text(s, val, In(7.75), In(1.73)+i*In(1.1), In(2.5), In(0.45),
             font_size=18, bold=True, color=color)
    add_text(s, lbl, In(10.3), In(1.73)+i*In(1.1), In(2.5), In(0.55),
             font_size=11, color=C_GRAY)

add_rect(s, In(7.6), In(6.15), In(5.4), In(0.6),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=C_RED, line_pt=0.8)
add_text(s, "Mesmo com educação, sexo, idade e local de moradia IDÊNTICOS, negros têm 25,3% menor odds de estar em ocupação qualificada (OR=0,747 GLMM M2).",
         In(7.75), In(6.2), In(5.1), In(0.55),
         font_size=12, bold=True, color=C_RED)
footer(s, 10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — REGRESSÃO QUANTÍLICA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "10. Regressão Quantílica — Formalização do Glass Ceiling",
           "O gap racial aumenta no topo da distribuição de renda — confirmação formal")

add_img(s, FIGURES / "quantreg_trajetoria.png", In(0.3), In(1.2), In(7.2))

add_text(s, "O que o gráfico mostra", In(7.8), In(1.2), In(5.3), In(0.4),
         font_size=16, bold=True, color=C_DARK)
bullet_box(s, [
    "Eixo Y: gap racial em % (quanto negros ganham a menos)",
    "Cada ponto: estimativa para aquele percentil da distribuição",
    "Inclinação descendente = gap CRESCE no topo = Glass Ceiling",
    "M3 vs M4: a diferença é a mediação ocupacional — maior no topo",
    "No q95, mais da metade do gap é explicada pela exclusão de ocupações de alta remuneração",
], In(7.8), In(1.7), In(5.3), In(3.0), font_size=13, dot_color=C_RED)

add_rect(s, In(7.8), In(5.0), In(5.3), In(1.7),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=C_RED, line_pt=1)
add_text(s, "Glass Ceiling Racial: confirmado",
         In(8.0), In(5.1), In(5.0), In(0.4),
         font_size=15, bold=True, color=C_RED)
add_text(s, "O gap não é uniforme — cresce sistematicamente quanto maior a remuneração. Isso significa que a barreira racial é mais intensa exatamente onde os prêmios de progressão são maiores.",
         In(8.0), In(5.5), In(4.9), In(1.1),
         font_size=12, color=C_BLACK)
footer(s, 11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — RIF-OB: STICKY FLOOR vs GLASS CEILING POR DOTAÇÕES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "11. RIF-OB — Sticky Floor Discriminatório vs. Glass Ceiling por Dotações",
           "Firpo, Fortin & Lemieux (2018) | N=7.694.198 | Decomposição incondicional por quantil")

add_img(s, FIGURES / "rif_ob_retornos_quantis.png", In(0.3), In(1.2), In(6.8))

add_text(s, "O que a RIF-OB revela além da QR?", In(7.5), In(1.2), In(5.6), In(0.4),
         font_size=15, bold=True, color=C_DARK)

for j, (hdr, xh, wc) in enumerate(zip(
        ["Quantil", "Dotações", "Retornos"],
        [In(7.5), In(9.9), In(11.45)],
        [In(2.38), In(1.53), In(1.53)])):
    add_rect(s, xh, In(1.72), wc, In(0.42), fill_rgb=C_DARK)
    add_text(s, hdr, xh+In(0.05), In(1.73), wc-In(0.1), In(0.38),
             font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

rif_table = [
    ("q10  (base)",    "61,3%", "33,1%", RGBColor(0xFF,0xEB,0xEE), C_RED),
    ("q50  (mediana)", "57,4%", "21,5%", C_LGRAY,                  C_GRAY),
    ("q90  (topo)",    "75,8%", "11,2%", RGBColor(0xE3,0xF2,0xFD), C_BLUE),
]
for i, (ql, d, r, bg, clr) in enumerate(rif_table):
    add_rect(s, In(7.5), In(2.18)+i*In(0.65), In(5.5), In(0.6),
             fill_rgb=bg, line_rgb=clr, line_pt=0.6)
    add_text(s, ql, In(7.55), In(2.22)+i*In(0.65), In(2.3), In(0.5),
             font_size=12, color=C_BLACK)
    add_text(s, d,  In(9.9),  In(2.22)+i*In(0.65), In(1.53), In(0.5),
             font_size=13, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, r,  In(11.45),In(2.22)+i*In(0.65), In(1.53), In(0.5),
             font_size=13, bold=True, color=C_RED,  align=PP_ALIGN.CENTER)

bullet_box(s, [
    "Retornos DECLINAM: q10=33,1% → q90=11,2%  (Δ = −21,9 pp)",
    "Dotações CRESCEM: q10=61,3% → q90=75,8%  (Δ = +14,4 pp)",
    "Glass ceiling salarial causado por desvantagens PRÉ-MERCADO\n(educação, redes, segregação residencial)",
    "Não por discriminação crescente no topo da hierarquia salarial",
], In(7.5), In(4.18), In(5.6), In(2.0), font_size=13, dot_color=C_RED)

add_rect(s, In(0.3), In(6.3), In(12.7), In(0.55),
         fill_rgb=RGBColor(0xFF,0xEB,0xEE), line_rgb=C_RED, line_pt=1)
add_text(s, "Sticky Floor: discriminação de mercado é 3× mais intensa na base (q10=33,1%) do que no topo (q90=11,2%). "
            "Enforcement antidiscriminação ataca o q10; reformar dotações (acesso, redes) erode o glass ceiling no topo.",
         In(0.5), In(6.35), In(12.3), In(0.45),
         font_size=12, bold=True, color=C_DARK)
footer(s, 12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — INTERSECCIONALIDADE OB 4 GRUPOS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "12. Interseccionalidade — OB em 4 Grupos (Crenshaw, 1989)",
           "OB twofold vs. Homem Branco | N=2.357.851 | penalidade interseccional extra = +9,5 pp")

add_img(s, FIGURES / "interseccional_ob4grupos.png", In(0.3), In(1.2), In(6.5))

add_text(s, "Resultados OB 4 Grupos", In(7.1), In(1.2), In(5.9), In(0.4),
         font_size=15, bold=True, color=C_DARK)

grupos_int = [
    (C_BLUE,                       "Mulher Branca — gap 46,6%",
     "Retornos dominam: 109%\nDiscriminação de gênero pura\n(dotações ligeiramente negativas: −8,9%)"),
    (C_RED,                        "Homem Negro — gap 40,3%",
     "Dotações dominam: 71%\nGap majoritariamente estrutural\n(capital humano acumulado inferior)"),
    (RGBColor(0x6A, 0x00, 0x8A),   "Mulher Negra — gap 96,4%",
     "Penalidade interseccional extra: +9,5 pp\nacima da soma Mulher Branca + Homem Negro\n(Crenshaw, 1989 — confirmado empiricamente)"),
]
for i, (color, titulo, desc) in enumerate(grupos_int):
    add_rect(s, In(7.1), In(1.72)+i*In(1.68), In(5.9), In(1.52),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=1.5)
    add_rect(s, In(7.1), In(1.72)+i*In(1.68), In(5.9), In(0.48), fill_rgb=color)
    add_text(s, titulo, In(7.22), In(1.75)+i*In(1.68), In(5.65), In(0.45),
             font_size=13, bold=True, color=C_WHITE)
    add_text(s, desc, In(7.22), In(2.28)+i*In(1.68), In(5.65), In(0.92),
             font_size=12, color=C_BLACK)

add_rect(s, In(0.3), In(6.78), In(12.7), In(0.56),
         fill_rgb=RGBColor(0xE3,0xF2,0xFD), line_rgb=C_DARK, line_pt=0.8)
add_text(s, "HLM interseccional: β_neg×fem = +0,069 (p<0,001) → penalidade combinada (−40,5%) < soma isolada (−44,5%): "
            "as duas desvantagens competem pelo mesmo mecanismo de exclusão — efeito de compressão.",
         In(0.5), In(6.83), In(12.3), In(0.45),
         font_size=12, color=C_DARK)
footer(s, 13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — SENSIBILIDADE A VARIÁVEIS OMITIDAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "13. Sensibilidade — Konfound, E-values e Oster Bounds",
           "3 métodos complementares — nenhum consegue anular o gap racial observado")

add_img(s, FIGURES / "sensibilidade_konfound_evalues.png", In(0.3), In(1.2), In(6.8))

add_text(s, "Resultados de robustez", In(7.5), In(1.2), In(5.6), In(0.4),
         font_size=15, bold=True, color=C_DARK)

rob_blocks = [
    (C_DARK,  "Konfound HLM (Frank et al., 2013)",
     "pkonfound: 99,5% (M1)  |  98,8% (M2/M3)  |  98,5% (M4)",
     "99,5% dos casos precisariam ser removidos para anular M1"),
    (C_AMBER, "OLS subestima robustez — Δ = 5,4–9,7 pp",
     "OLS: 91,5% / 93,5% / 88,7%  →  HLM: +7,9 / +5,4 / +9,7 pp",
     "SEs OLS são 6–16× maiores: não captura clustering → robustez subestimada"),
    (C_RED,   "E-values GLMM (VanderWeele & Ding, 2017)",
     "E ≥ 2,17× (M1)  |  E ≥ 2,04× (M2)  para anular OR = 0,704 / 0,747",
     "Confounder precisaria ter associação ≥ 2× com raça E com ocupação — implausível"),
    (C_BLUE,  "Oster Bounds OLS auxiliar (Oster, 2019)",
     "δ* ∈ {−0,48; −0,43; −0,39} — todos negativos",
     "Omitidas precisariam agir na direção OPOSTA aos observáveis — impossível"),
]
for i, (color, titulo, resultado, interp) in enumerate(rob_blocks):
    add_rect(s, In(7.5), In(1.72)+i*In(1.25), In(5.6), In(1.15),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=1.2)
    add_text(s, titulo, In(7.65), In(1.76)+i*In(1.25), In(5.3), In(0.35),
             font_size=11.5, bold=True, color=color)
    add_text(s, resultado, In(7.65), In(2.12)+i*In(1.25), In(5.3), In(0.38),
             font_size=12.5, bold=True, color=C_BLACK)
    add_text(s, interp, In(7.65), In(2.5)+i*In(1.25), In(5.3), In(0.3),
             font_size=10.5, italic=True, color=C_GRAY)

footer(s, 14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — TENDÊNCIA TEMPORAL + HECKMAN + EVENT STUDY COVID
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "14. Análises Temporais e Seleção — Chow, Heckman e Event Study COVID",
           "Quebra estrutural 2020 | Seleção corrigida (Heckman) | COVID: convergência aparente de curto prazo")

add_img(s, FIGURES / "tendencia_temporal.png",  In(0.3), In(1.2), In(4.1), In(3.5))
add_img(s, FIGURES / "heckman_selecao.png",     In(4.6), In(1.2), In(4.1), In(3.5))
add_img(s, FIGURES / "event_study_covid.png",   In(9.1), In(1.2), In(4.0), In(3.5))

cap_blocks = [
    (C_AMBER, "Chow test — Quebra 2020",
     "F(2,6) = 7,012   p = 0,027\nQuebra estrutural detectada em 2020\n"
     "Tendência geral: δ = 0,0008 log-pt/ano (estável, p=0,077)\n"
     "Gap racial não convergiu espontaneamente em 10 anos"),
    (C_RED,   "Heckman — Correção de Seleção",
     "OLS: gap = −9,71%\nHeckman (com IMR): gap = −7,41%\n"
     "Δ = +2,3 pp  |  λ = −1,985  (p < 0,001)\n"
     "OLS superestima gap: negros empregados aceitam\nsalários inferiores (seleção por tolerância)"),
    (C_BLUE,  "Event Study COVID (DiD)",
     "τ = +0,015  (SE = 0,007,  p = 0,025)\nNegros sofreram MENOS com COVID\n"
     "no curto prazo — convergência aparente\n"
     "DiD corrigido: 0,015 vs. ingênuo: 0,023"),
]
for i, (color, titulo, corpo) in enumerate(cap_blocks):
    x = In(0.3) + i * In(4.35)
    add_rect(s, x, In(4.85), In(4.1), In(2.3),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=1.5)
    add_rect(s, x, In(4.85), In(4.1), In(0.48), fill_rgb=color)
    add_text(s, titulo, x+In(0.12), In(4.88), In(3.86), In(0.45),
             font_size=12, bold=True, color=C_WHITE)
    add_text(s, corpo, x+In(0.12), In(5.42), In(3.86), In(1.65),
             font_size=11, color=C_BLACK)

footer(s, 15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ML/SHAP
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "15. ML/SHAP — O Que Mais Determina a Renda?",
           "XGBoost R²=0,6162 | Horas + CBO + Formalidade emergem como top preditores")

add_img(s, FIGURES / "shap_importance_xgb.png", In(0.3), In(1.2), In(6.8))

add_text(s, "Destaques com as novas variáveis", In(7.4), In(1.2), In(5.8), In(0.4),
         font_size=16, bold=True, color=C_DARK)

ranking = [
    ("#1", "Renda média da UPA", "0,272", "Wilson (1987): onde mora supera o quanto estudou"),
    ("#2", "Horas trabalhadas", "0,166", "Novo — VD4031 extraído do ZIP bruto"),
    ("#3", "CBO: Profissionais", "0,119", "Novo — V4010 primeiro dígito"),
    ("#4", "Emprego formal", "0,109", "Novo — VD4009 carteira/público"),
    ("#11","Raça (negro)", "0,029", "−2,5% residual pós-ocupação"),
]
for i, (rank, feat, shap, note) in enumerate(ranking):
    color = C_RED if rank == "#11" else (C_BLUE if rank == "#1" else C_BLACK)
    add_rect(s, In(7.4), In(1.7)+i*In(0.88), In(5.8), In(0.82),
             fill_rgb=RGBColor(0xFF,0xEB,0xEE) if rank == "#11" else C_LGRAY,
             line_rgb=C_RED if rank == "#11" else C_GRAY, line_pt=0.5)
    add_text(s, rank, In(7.5), In(1.75)+i*In(0.88), In(0.5), In(0.4),
             font_size=12, bold=True, color=color)
    add_text(s, feat, In(8.05), In(1.75)+i*In(0.88), In(2.5), In(0.4),
             font_size=12, bold=(rank in ["#1","#11"]), color=color)
    add_text(s, f"|SHAP|={shap}", In(10.65), In(1.75)+i*In(0.88), In(1.0), In(0.4),
             font_size=11, bold=True, color=color)
    add_text(s, note, In(8.05), In(2.1)+i*In(0.88), In(4.9), In(0.3),
             font_size=9.5, color=C_GRAY, italic=True)

footer(s, 16)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SNA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "16. Análise de Redes Sociais — Capital Social e Mobilidade",
           "Betweenness centrality nula para negros em todos os níveis educacionais")

add_img(s, FIGURES / "sna_rede_demografica.png", In(0.3), In(1.2), In(6.5))

add_text(s, "O que é betweenness centrality?", In(7.2), In(1.2), In(5.9), In(0.4),
         font_size=16, bold=True, color=C_DARK)
bullet_box(s, [
    "Mede quem 'conecta' grupos distintos na rede",
    "Alta betweenness = broker = acesso antecipado a vagas, mentoria, promoções (Burt, 2004)",
    "BRANCOS: betweenness elevada em todos os níveis educacionais",
    "NEGROS: betweenness = 0 mesmo com pós-graduação",
    "Implicação: credenciais educacionais de negros valem menos porque falta acesso às redes de conversão",
], In(7.2), In(1.7), In(5.9), In(3.5), font_size=13, dot_color=C_DARK)

add_rect(s, In(7.2), In(5.45), In(5.9), In(0.85),
         fill_rgb=RGBColor(0xE3,0xF2,0xFD), line_rgb=C_DARK, line_pt=0.8)
add_text(s, '"Negros enfrentam um duplo obstáculo ao retorno educacional: o gap direto mensurado pelo HLM e a exclusão das redes que convertem diplomas em empregos."',
         In(7.35), In(5.5), In(5.6), In(0.8),
         font_size=12, italic=True, color=C_DARK)
footer(s, 17)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ESTADO E DESIGUALDADE (H1 / H2 / H3)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "17. Estado e Desigualdade — H1 / H2 / H3",
           "Setor público reduz desigualdade geral mas gap racial persiste — e gap de gênero é MAIOR no público")

add_img(s, FIGURES / "estado_h1_gini.png",   In(0.3), In(1.2), In(6.3))
add_img(s, FIGURES / "estado_h2h3_gaps.png", In(6.8), In(1.2), In(6.2))

kpi_box(s, "Gini setor público",    "0,466",   "vs privado 0,472 | total 0,488",      In(0.3),  In(4.6), w=In(2.8), val_color=C_BLUE)
kpi_box(s, "Gap racial público",    "−25,2%",  "controlado | privado −29,2%",         In(3.3),  In(4.6), w=In(3.1), val_color=C_RED)
kpi_box(s, "Gap gênero público",    "−18,8%",  "PIOR que privado −16,9% (paradoxo)",  In(6.6),  In(4.6), w=In(3.5), val_color=C_AMBER)
kpi_box(s, "Prêmio salarial pub.",  "+99,6%",  "Theil entre setores: 7,8% do total",  In(10.3), In(4.6), w=In(2.7), val_color=C_GREEN)

add_rect(s, In(0.3), In(6.05), In(12.7), In(0.50),
         fill_rgb=RGBColor(0xE3,0xF2,0xFD), line_rgb=C_DARK, line_pt=0.8)
add_text(s, "H1: Estado reduz levemente o Gini (0,466 vs 0,472) mas prêmio de +99,6% concentra renda entre servidores.  "
            "H2: Racismo persiste no público (−25,2%) — atenuado vs privado (−29,2%).  "
            "H3 (paradoxo): concurso iguala entrada, mas promoção/DAS favorece homens → gap de gênero maior no público.",
         In(0.5), In(6.10), In(12.3), In(0.45),
         font_size=11.5, color=C_DARK)

footer(s, 18)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — RENDA REAL, ARMADILHA E INCLUSÃO (H4 / H5)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "18. Renda Real, Armadilha e Inclusão — H4 / H5",
           "Pleno emprego sem prosperidade + gap de qualificação quase inalterado em 10 anos")

add_img(s, FIGURES / "estado_h4_tendencia.png", In(0.3), In(1.2), In(6.3))
add_img(s, FIGURES / "estado_h5_armadilha.png", In(6.8), In(1.2), In(6.2))

kpi_box(s, "Renda real 2025",     "R$1.283",  "vs pico R$1.345 em 2024",              In(0.3),  In(4.6), w=In(3.0), val_color=C_RED)
kpi_box(s, "Emprego 2025",        "94,4%",    "máxima histórica da série (2016–2025)", In(3.5),  In(4.6), w=In(2.9), val_color=C_GREEN)
kpi_box(s, "Alta qualif. negros", "11,0%",    "brancos: 22,4% — gap estável +11pp",   In(6.6),  In(4.6), w=In(3.5), val_color=C_RED)
kpi_box(s, "Ganho de inclusão",   "+127,3%",  "renda negra se CBO = branco",          In(10.3), In(4.6), w=In(2.7), val_color=C_GREEN)

add_rect(s, In(0.3), In(6.05), In(12.7), In(0.50),
         fill_rgb=RGBColor(0xFF,0xF9,0xE7), line_rgb=C_AMBER, line_pt=1)
add_text(s, "H4: Contradição do mercado atual — emprego em máxima histórica (94,4%) mas renda real recua em 2025 (−4,6% vs 2024).  "
            "H5: Armadilha da renda média — gap de qualificação (11,0% vs 22,4%) praticamente inalterado após 10 anos. "
            "Inclusão produtiva plena = proxy +74,3 p.p. de PIB (Hsieh et al., 2019).",
         In(0.5), In(6.10), In(12.3), In(0.45),
         font_size=11.5, color=C_DARK)

footer(s, 19)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — JUSTIFICAÇÃO METODOLÓGICA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "19. Por Que Esses Modelos? — Justificação Estatística",
           "LRT χ²=191.625 confirma hierarquia | ICC=9,83% > limiar 5% | HLM supera OLS+FE em AIC")

add_img(s, FIGURES / "modelos_loglik_aic.png", In(0.3), In(1.2), In(6.3))
add_img(s, FIGURES / "modelos_lrt_icc.png",   In(6.8), In(1.2), In(6.2))

kpi_box(s, "LRT HLM vs OLS (nulo)", "191.625",   "χ² Δk=1 — hierarquia confirmada",     In(0.3),  In(4.55), w=In(3.9), val_color=C_DARK)
kpi_box(s, "ICC por UPA",           "9,83%",     "> 5% — Raudenbush & Bryk (2002)",     In(4.4),  In(4.55), w=In(3.5), val_color=C_BLUE)
kpi_box(s, "AIC HLM Contextual",    "3.588.684", "OLS+FE(UF): 3.684.833 (+96k pior)",   In(8.1),  In(4.55), w=In(5.0), val_color=C_GREEN)

bullet_box(s, [
    "OLS ignora correlação intra-UPA → erros padrão subestimados → inferência inválida (θ OLS ≠ θ BLUP)",
    "OLS+FE(UF) piora AIC vs OLS Individual: dummies de UF são grosseiras demais para capturar variação intra-estado",
    "HLM 3 níveis (indivíduo > UPA > UF) é o único estimador com estrutura compatível com os dados e inferência válida",
], In(0.3), In(5.95), In(12.7), In(1.15), font_size=13, dot_color=C_DARK)

footer(s, 20)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SÍNTESE: TRIÂNGULO DE EVIDÊNCIAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "20. Síntese — Triângulo de Evidências",
           "5 métodos independentes apontam para o mesmo diagnóstico")

# Três vértices do triângulo
vertices = [
    (In(0.3),  In(1.3),  C_RED,   "DISCRIMINAÇÃO\nDE ACESSO",
     ["GLMM: OR=0,747 para ocp. qualificada", "AME=−1,07 p.p. residual (M2)", "Persiste após todos os controles observáveis"]),
    (In(6.9),  In(1.3),  C_BLUE,  "DISCRIMINAÇÃO\nDE REMUNERAÇÃO",
     ["HLM M4: gap residual 6,2%", "Quantile Reg.: cresce no topo (glass ceiling)", "SHAP: −2,5% efeito racial residual"]),
    (In(3.6),  In(4.3),  C_DARK,  "EXCLUSÃO\nDAS REDES",
     ["SNA: betweenness=0 para negros", "Diplomas valem menos sem acesso às redes", "Cluster: dupla desvantagem mulher negra"]),
]
for l, t, color, title, items in vertices:
    add_rect(s, l, t, In(5.8), In(2.6),
             fill_rgb=RGBColor(0xF5,0xF5,0xF5), line_rgb=color, line_pt=1.5)
    add_rect(s, l, t, In(5.8), In(0.5), fill_rgb=color)
    add_text(s, title, l+In(0.1), t+In(0.05), In(5.6), In(0.45),
             font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for i, item in enumerate(items):
        add_text(s, "• " + item, l+In(0.2), t+In(0.6)+i*In(0.6), In(5.4), In(0.55),
                 font_size=12, color=C_BLACK)

add_rect(s, In(0.3), In(6.4), In(12.7), In(0.5),
         fill_rgb=RGBColor(0x1F,0x38,0x64), line_rgb=C_AMBER, line_pt=0)
add_text(s, "Oaxaca-Blinder une os três vértices: 84% do gap é de ACESSO (dotações) | 16% de REMUNERAÇÃO (retornos) | SNA explica por que o acesso é estruturalmente negado  |  GLMM: OR=0,747",
         In(0.5), In(6.45), In(12.3), In(0.45),
         font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
footer(s, 21)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — PO: PESQUISA OPERACIONAL — PRIORIZAÇÃO DE POLÍTICAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "21. Pesquisa Operacional — Priorização Multicritério Anti-Discriminação",
           "TOPSIS + AHP (CR=0,004) + Programação Linear + Fronteira de Pareto | 6 políticas avaliadas")

add_img(s, FIGURES / "po_politicas_topsis.png", In(0.3), In(1.2), In(7.2))

add_text(s, "Ranking TOPSIS (6 intervenções)", In(8.0), In(1.2), In(5.1), In(0.4),
         font_size=14, bold=True, color=C_DARK)

topsis_rows = [
    (1, "Cotas ocupacionais CBO 1–4",        "CC = 0,799", C_RED),
    (2, "Equidade educacional",               "CC = 0,558", C_BLUE),
    (3, "Mentoria e redes profissionais",     "CC = 0,388", C_GREEN),
    (4, "Transparência salarial obrigatória", "CC = 0,324", C_AMBER),
    (5, "Enforcement anti-discriminação",     "CC = 0,241", C_GRAY),
    (6, "Desegregação residencial",           "CC = 0,109", C_GRAY),
]
for i, (rank, nome, cc, color) in enumerate(topsis_rows):
    bg = RGBColor(0xFF,0xEB,0xEE) if rank == 1 else (RGBColor(0xE3,0xF2,0xFD) if rank == 2 else C_LGRAY)
    add_rect(s, In(8.0), In(1.72)+i*In(0.8), In(5.1), In(0.74),
             fill_rgb=bg, line_rgb=color, line_pt=0.8)
    add_text(s, f"#{rank}", In(8.1), In(1.77)+i*In(0.8), In(0.4), In(0.6),
             font_size=13, bold=True, color=color)
    add_text(s, nome, In(8.55), In(1.77)+i*In(0.8), In(3.2), In(0.6),
             font_size=12, bold=(rank <= 3), color=C_BLACK)
    add_text(s, cc, In(11.8), In(1.77)+i*In(0.8), In(1.25), In(0.6),
             font_size=12, bold=True, color=color, align=PP_ALIGN.RIGHT)

add_rect(s, In(0.3), In(6.6), In(12.7), In(0.65),
         fill_rgb=RGBColor(0x1F,0x38,0x64), line_rgb=C_AMBER, line_pt=0)
add_text(s, "84% do gap é de ACESSO → P1 (Cotas CBO) lidera com CC=0,799 — 2× superior ao 2º colocado.  "
            "Pareto (λ=0,5): alocação ótima = P1+P5 (cotas + mentoria).  "
            "PL-1 (orçamento B=5): redução projetada de 24% do gap bruto.",
         In(0.5), In(6.65), In(12.3), In(0.55),
         font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
footer(s, 22)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — IMPLICAÇÕES DE POLÍTICA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "22. Implicações de Política",
           "Três eixos simultâneos — ações isoladas são insuficientes para romper a armadilha estrutural")

politicas = [
    (C_RED,   "Eixo 1 — Acesso",
     ["Cotas de acesso CBO Dirigente/Profissional:\nmeta IR ≥ 0,80 até 2030",
      "PRONATEC c/ recorte racial e subsídio de\ntransporte para trabalhadores de UPAs pobres",
      "Bolsas de residência profissional em empresas\nde alta remuneração para egressos negros"]),
    (C_BLUE,  "Eixo 2 — Remuneração",
     ["Transparência salarial por raça/gênero\nobrigatória (empresas > 100 funcionários)",
      "Auditoria de igual pagamento por trabalho igual\ncom penalidades progressivas (gap HLM M4)",
      "Piso salarial indexado nas categorias com\nmaior gap racial residual (6,2% HLM M4)"]),
    (C_GREEN, "Eixo 3 — Inclusão Produtiva",
     ["Mentoria estruturada para elevar betweenness\nde negros nas redes profissionais (SNA: =0)",
      "30% dos cargos DAS e liderança corporativa\npara negros até 2030 (inclusão nas redes)",
      "Equalizar CBO = +127,3% renda negra |\nproxy PIB +74,3 p.p. (Hsieh et al., 2019)"]),
]
for i, (color, title, items) in enumerate(politicas):
    x = In(0.3) + i * In(4.35)
    add_rect(s, x, In(1.2), In(4.1), In(5.5),
             fill_rgb=C_LGRAY, line_rgb=color, line_pt=1.5)
    add_rect(s, x, In(1.2), In(4.1), In(0.55), fill_rgb=color)
    add_text(s, title, x+In(0.1), In(1.25), In(3.9), In(0.5),
             font_size=13, bold=True, color=C_WHITE)
    for j, item in enumerate(items):
        add_text(s, "▸  " + item, x+In(0.2), In(1.85)+j*In(0.85), In(3.7), In(0.8),
                 font_size=12, color=C_BLACK)

add_text(s, "⚠  Ao ritmo atual de convergência (~0,02 log-pontos/ano), eliminar o gap levaria mais de 100 anos.",
         In(0.3), In(6.8), In(12.7), In(0.45),
         font_size=13, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
footer(s, 23)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — LIMITAÇÕES E AGENDA FUTURA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header_bar(s, "23. Limitações e Agenda Futura",
           "Honestidade acadêmica: o que este trabalho não faz e por quê")

add_text(s, "Limitações", In(0.4), In(1.25), In(6.0), In(0.4),
         font_size=16, bold=True, color=C_RED)
bullet_box(s, [
    "PNAD não permite experimentos causais — coeficientes são associações condicionais, não efeitos causais no sentido de Rubin/Pearl",
    "UPA como proxy de bairro: unidade de amostragem ≠ bairro administrativo",
    "GLMM estimado na população completa via lme4::glmer (nAGQ=0, bobyqa, n=2,4M, 39,8k UPAs) — M3 com inclinação aleatória de negro não estimado por custo computacional",
    "CBO auto-declarado pode ter viés de classificação por raça (deflation de ocupação)",
], In(0.4), In(1.75), In(6.2), In(3.8), font_size=13, dot_color=C_RED)

add_text(s, "Agenda Futura", In(7.0), In(1.25), In(6.0), In(0.4),
         font_size=16, bold=True, color=C_BLUE)
bullet_box(s, [
    "Análise de variáveis instrumentais (IV) para identificação causal — distância à escola técnica como instrumento",
    "RAIS linkada à PNAD para rastrear trajetórias de mobilidade intraempresa",
    "Análise de mediação formal (path analysis) raça → ocupação → renda",
    "Extensão ao setor público por UF — heterogeneidade do glass ceiling entre estados",
    "Modelo longitudinal com painéis rotativos da PNAD (2T seguidos) para efeitos fixos de indivíduo",
], In(7.0), In(1.75), In(6.2), In(3.8), font_size=13, dot_color=C_BLUE)

footer(s, 24)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CONCLUSÃO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill_rgb=C_DARK)
add_rect(s, 0, 0, W, In(0.8), fill_rgb=RGBColor(0x0D,0x1F,0x3C))
add_text(s, "18. Conclusão", In(0.4), In(0.1), In(12), In(0.65),
         font_size=26, bold=True, color=C_WHITE, font_name="Calibri")

numeros = [
    ("19,3%",    "Gap racial bruto\n(M1 HLM)"),
    ("6,2%",     "Discriminação pura\n(M4 HLM)"),
    ("OR=0,747", "Acesso qualificado\n(GLMM M2, lme4)"),
    ("χ²=191k",  "LRT confirma\nhierarquia UPA"),
]
for i, (val, lbl) in enumerate(numeros):
    x = In(0.4) + i * In(3.2)
    add_rect(s, x, In(0.95), In(3.0), In(1.5),
             fill_rgb=RGBColor(0x1A,0x30,0x5C), line_rgb=C_AMBER, line_pt=1)
    add_text(s, val, x+In(0.1), In(1.0), In(2.8), In(0.75),
             font_size=30, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x+In(0.1), In(1.75), In(2.8), In(0.6),
             font_size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

conclusoes = [
    "A desigualdade racial no mercado de trabalho brasileiro é estrutural, multicausal e resistente à convergência espontânea.",
    "A discriminação opera em dois estágios independentes: barreiras de ACESSO a ocupações qualificadas (GLMM: OR=0,747, AME=−1,07 p.p.) e discriminação de REMUNERAÇÃO dentro das mesmas funções (HLM M4: 6,2%).",
    "O glass ceiling racial é real e crescente no topo da distribuição — confirmado formalmente pela regressão quantílica e consistente com o IR=0,47 nas capitais.",
    "Políticas que atuam apenas no gap salarial direto atacam 16% do problema. Os 84% restantes requerem ação nas portas de entrada das ocupações de alto prestígio.",
]
for i, texto in enumerate(conclusoes):
    add_rect(s, In(0.3), In(2.6)+i*In(0.95), In(12.7), In(0.82),
             fill_rgb=RGBColor(0x0F,0x27,0x4A), line_rgb=C_AMBER if i==0 else C_BLUE, line_pt=0.6)
    add_text(s, f"{i+1}. " + texto, In(0.5), In(2.65)+i*In(0.95), In(12.3), In(0.75),
             font_size=13, color=C_WHITE, font_name="Calibri")

add_rect(s, 0, H-In(0.5), W, In(0.5), fill_rgb=RGBColor(0x0D,0x1F,0x3C))
add_text(s, "Ricardo Calheiros  |  MBA Data Science & Analytics  |  USP/ESALQ  |  rickinrj@gmail.com",
         In(0.3), H-In(0.45), In(12.7), In(0.4),
         font_size=11, color=RGBColor(0x90,0xA4,0xAE), align=PP_ALIGN.CENTER)

# ── Salvar ────────────────────────────────────────────────────────────────────
prs.save(str(OUT_PPT))
print(f"Arquivo gerado: {OUT_PPT.name}")
print(f"Tamanho: {OUT_PPT.stat().st_size // 1024} KB")
print("Abra com PowerPoint para revisar e ajustar o layout.")
